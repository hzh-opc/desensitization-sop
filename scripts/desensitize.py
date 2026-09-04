#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
一键脱敏本地脚本（desensitize.py）

功能：在本地对文本 / 数据文件中的敏感标识符进行脱敏，生成脱敏副本，
并将「可逆映射表」加密保存到独立目录，确保原始文件不离开本机。

设计原则（对齐《敏感信息脱敏操作 SOP》）：
- 本地优先：脚本只在本地运行，原始文件与映射表永不离开本机。
- 最小必要：只处理识别到的敏感字段。
- 映射表分离 + AES 加密：脱敏副本可上云；映射表单独存放并加密（密钥不入库）。
- 中文优先：内置中文身份证 / 手机号 / 银行卡 / 邮箱 / 车牌 / 护照等正则。
- 中文增强（--cn-enhance）：本地正则识别中文姓名 / 地址 / 机构名，
  纯离线、无需安装额外模型。

依赖：
- 必需：cryptography（AES / Fernet）
- 解密：pikepdf（封装 QPDF，稳健解密 PDF）、msoffcrypto-tool（解密 Office）
- PDF 文本抽取 / 页面渲染：pypdfium2（Chrome 同款 PDFium，C 实现，快且稳健）
- 本地 OCR：rapidocr + onnxruntime（纯 pip 安装，模型随 wheel 捆绑，
  完全离线、数据不出本机；中文识别基于 PP-OCRv6 模型）
- 中文增强：纯内置正则，无额外依赖（不依赖 Presidio 等需联网下载模型的方案）

Python 版本：标准（非 free-threaded）CPython >=3.10 且 <3.14——
onnxruntime 目前不提供 free-threaded wheel，3.14 支持亦未完备，
故 Python 版本按全部依赖的兼容性综合确定。

用法：
  # 扫描（仅报告命中，不生成文件）
  uv run --project <tool_dir> python desensitize.py scan ./data/

  # 脱敏（默认 hybrid 模式），输出到 ./desensitized，映射表到 ./.desensitize_keys
  uv run --project <tool_dir> python desensitize.py run ./data/ \
      --out ./desensitized --keys ./.desensitize_keys

  # 令牌化（可逆，便于跨记录关联计数 / 分组）
  uv run --project <tool_dir> python desensitize.py run ./data/ --mode token

  # 语义掩码 + 唯一令牌（hybrid）：保留字段语义（张*/138****8000…）且可无歧义恢复
  uv run --project <tool_dir> python desensitize.py run ./data/ --mode hybrid --cn-enhance

  # 中文识别增强：额外识别中文姓名 / 地址 / 机构名
  uv run --project <tool_dir> python desensitize.py run ./data/ --cn-enhance
  uv run --project <tool_dir> python desensitize.py scan ./data/ --cn-enhance --recursive

  # 用口令派生密钥（密钥不落盘，靠口令恢复映射表）
  uv run --project <tool_dir> python desensitize.py run ./data/ --passphrase "***"

  # 本地复核：解密映射表，验证可逆性（原始值不出本机、无需上云）
  uv run --project <tool_dir> python desensitize.py decrypt --keys ./.desensitize_keys
  uv run --project <tool_dir> python desensitize.py decrypt --keys ./.desensitize_keys --out mapping_review.json

  # 回填：用映射表把脱敏副本还原为含原值的本地内部文档（映射表不离本地）
  uv run --project <tool_dir> python desensitize.py restore --keys ./.desensitize_keys \
      --input ./desensitized --out ./restored
  # 仅回填姓名与身份证（如工资条只需回写姓名）
  uv run --project <tool_dir> python desensitize.py restore --keys ./.desensitize_keys \
      --types name,id_card --out ./restored_names

  # 审计：基于 run 报告自动生成 12 项审计文档
  uv run --project <tool_dir> python desensitize.py audit \
      --report ./desensitize_report.json --out ./desensitize_audit.md

  # 本地预处理（v2.4 起自动解密 + 纯本地 rapidocr OCR）
  uv run --project <tool_dir> python desensitize.py preprocess ./data/ \
      --passwords-file ./pw.txt --out-dir ./preprocessed
  # 确认无误并校对 OCR 后，run 会再次校验无异常才放行
  uv run --project <tool_dir> python desensitize.py run ./preprocessed \
      --preprocess-manifest ./preprocessed/desensitize_preprocess.json \
      --out ./desensitized --mode hybrid
"""

import argparse
import fnmatch
import glob
import hashlib
import json
import os
import re
import secrets
import shutil
import sys
from datetime import datetime

from cryptography.fernet import Fernet
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC

# 让本脚本无论以何种方式被加载（直接运行 / importlib 加载），都能找到同目录的
# desen_rules 辅助模块（识别规则与常量已拆至该模块）。
_SELF_DIR = os.path.dirname(os.path.abspath(__file__))
if _SELF_DIR not in sys.path:
    sys.path.insert(0, _SELF_DIR)

from desen_rules import (
    PATTERNS, build_patterns, _normalize_fullwidth, _KIND_VALIDATOR,
    SECRET_PATTERN, _sql_insert_secret_values,
    SECRET_EXTS, CODE_EXTS, TEXT_EXTS, EXTRACT_EXTS, IMAGE_EXTS,
    REMIND, WARN_ENCRYPTED, WARN_IMAGE, WARN_GENERIC,
    _find_cleaning_advice, CLEANING_ADVICE_TEXT, META_SKIP,
    _collect_tabular_names, _collect_tabular_names_xlsx,
)

# ---------------------------------------------------------------------------
# 统一成果中心（工作区，opt-in：仅当命令带 --workspace 才启用）
# 目的：把分散在各阶段的产物归拢到一处、用清晰中文目录命名、附“可上云/保密”
#       标签索引与上云前自检报告，方便非专业人员一站式查阅与核对。
# 不带 --workspace 时，各命令保持原默认行为（向后兼容，不影响 install 自检）。
# ---------------------------------------------------------------------------
WS_README = "00_使用说明.md"
WS_PREPROC_SUMMARY = "01_预处理确认单.md"
WS_UNDESEN = "02_未脱敏副本"
WS_READY = "02_未脱敏副本/解密与原始副本"
WS_OCR = "02_未脱敏副本/OCR待校对"
WS_DESEN = "03_脱敏副本"
WS_KEYS = "04_映射表_保密"
WS_SELFCHECK = "05_上云前自检报告.md"
WS_AUDIT = "06_审计与回填/审计记录.md"
WS_RESTORE = "06_审计与回填/回填成果"
WS_LOCAL_ONLY = "07_本地处理不外发"   # 🚫 不外发：仅本地处理、未脱敏原文，严禁上云、不得与脱敏副本一起存放
WS_INDEX_JSON = "成果索引.json"
WS_INDEX_MD = "成果索引.md"
# 映射表（重识别钥匙）是唯一含原值的文件，其路径本身即是攻击面：对话内与成果索引回显一律不展示其绝对路径。
WS_KEYS_MASK = "🔒 已加密封存（留本地，不展示路径）"


def _ws_root(args):
    """返回工作区绝对路径；未指定 --workspace 则返回 None。"""
    ws = getattr(args, "workspace", None)
    return os.path.abspath(ws) if ws else None


def _ensure_workspace(ws):
    """确保工作区目录存在，并在首次创建时写入使用说明。"""
    os.makedirs(ws, exist_ok=True)
    readme = os.path.join(ws, WS_README)
    if not os.path.exists(readme):
        _write_workspace_readme(readme, ws)
    return ws


def _write_workspace_readme(path, ws):
    L = [
        "# 脱敏工作区 · 使用说明",
        "",
        "本目录由「信息脱敏上云 SOP」自动整理，汇聚本次任务全部阶段性成果，方便你一站式查阅。",
        "",
        "## 目录含义与「能否上云」",
        "",
        "| 目录 / 文件 | 能否上云 | 含义 |",
        "| --- | --- | --- |",
        "| `03_脱敏副本/` | ✅ **可上云** | 已脱敏、可送云端模型处理的副本（唯一允许上传的目录） |",
        "| `02_未脱敏副本/` | 🚫 保密 | 原始 / 解密 / OCR 文本等未脱敏材料，**严禁外传** |",
        "| `04_映射表_保密/` | 🚫 保密 | 加密映射表（重识别钥匙），**绝不随副本上传** |",
        "| `06_审计与回填/回填成果/` | 🚫 保密 | 回填后的含原值内部文档，勿随副本外传 |",
        "| `07_本地处理不外发/` | 🚫 不外发 | 声明仅本地处理、未脱敏的原文，**严禁上云、不得与脱敏副本一起存放** |",
        "| `01_预处理确认单.md` | 内部 | 预处理确认：外发清单、OCR 校对、异常、确认闸门 |",
        "| `05_上云前自检报告.md` | 内部 | **上云前必读**：四处校对提醒 + 确认闸门 |",
        "| `06_审计与回填/审计记录.md` | 内部 | 对齐 12 项自查清单的审计文档 |",
        "| `成果索引.json` / `成果索引.md` | 内部 | 本索引，随时查看全部产物位置 |",
        "",
        "## 上云前必做（按顺序）",
        "1. 阅读 `05_上云前自检报告.md`；",
        "2. 校对 `02_未脱敏副本/OCR待校对/` 的 OCR 文本（逐字）；",
        "3. 抽查 `03_脱敏副本/` 脱敏是否到位、是否误伤业务；",
        "4. 确认无误后，由 AI Agent 在你明确「确认上云」后才上传 `03_脱敏副本/`。",
        "",
        "> 红线：原始敏感文件、未脱敏副本、映射表永远留本地，绝不整份上传。",
        "> 自动化识别非 100%（中文尤弱），禁止「一键脱敏即上云」，必须人工复核。",
    ]
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")


def _refresh_index(ws):
    """扫描工作区已知产物，生成 成果索引.json + 成果索引.md（幂等、自愈）。"""
    def _exists(rel):
        return os.path.exists(os.path.join(ws, rel))
    entries = []
    def _add(name, rel, tag, desc, hide_path=False):
        if _exists(rel):
            # 映射表（重识别钥匙）路径为攻击面：对话内 / 成果索引.md / status 回显一律以掩码替代。
            # 真实绝对路径仅写入 成果索引.json（本地文件，供本机复核，不默认回显到对话）。
            real = os.path.abspath(os.path.join(ws, rel))
            entries.append({"name": name, "path": real, "tag": tag, "desc": desc,
                            "display": WS_KEYS_MASK if hide_path else real})
    _add("预处理确认单", WS_PREPROC_SUMMARY, "内部", "外发清单、OCR 校对、异常清单、确认闸门")
    _add("未脱敏副本·解密与原始", WS_READY, "保密", "严禁上云：已解密未脱敏副本、原可直接处理文件")
    _add("未脱敏副本·OCR 待校对", WS_OCR, "保密", "严禁上云：OCR 识别文本，须逐字校对")
    _add("脱敏副本", WS_DESEN, "可上云", "★唯一允许上传云端的目录")
    _add("映射表", WS_KEYS, "保密", "重识别钥匙，绝不随副本上传", hide_path=True)
    _add("上云前自检报告", WS_SELFCHECK, "内部", "四处校对提醒 + 确认闸门，上云前必读")
    _add("审计记录", WS_AUDIT, "内部", "对齐 12 项自查清单的审计文档")
    _add("回填成果", WS_RESTORE, "保密", "含原值内部文档，勿随副本外传")
    _add("本地处理不外发", WS_LOCAL_ONLY, "不外发", "仅本地处理、未脱敏原文，严禁上云、不得与脱敏副本一起存放")
    idx = {"workspace": os.path.abspath(ws),
           "generated": datetime.now().isoformat(timespec="seconds"),
           "entries": entries}
    with open(os.path.join(ws, WS_INDEX_JSON), "w", encoding="utf-8") as f:
        json.dump(idx, f, ensure_ascii=False, indent=2)
    L = ["# 脱敏工作区 · 成果索引", "",
         "- 工作区：`%s`" % os.path.abspath(ws),
         "- 生成时间：%s" % idx["generated"], "",
         "## 各阶段成果（按可上云状态标注）", "",
         "| 阶段 / 产物 | 可上云状态 | 路径 | 说明 |",
         "| --- | --- | --- | --- |"]
    for e in entries:
        L.append("| %s | **%s** | %s | %s |" % (e["name"], e["tag"], e["display"], e["desc"]))
    L.append("")
    L.append("> 仅 `03_脱敏副本/` 可上云；`02_未脱敏副本/`、`04_映射表_保密/`、`06_审计与回填/回填成果/`、`07_本地处理不外发/` 均须留本地。")
    L.append("> 上云前请先阅读 `05_上云前自检报告.md` 并完成四处校对。")
    with open(os.path.join(ws, WS_INDEX_MD), "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")


def _write_selfcheck_report(ws, rep, ocr_dir):
    """生成《上云前自检报告》：汇总 OCR 副本 / 脱敏副本 / 本地不外发留痕 / 公开豁免留痕，
    并设确认闸门。"""
    out = os.path.join(ws, WS_SELFCHECK)
    counts = rep.get("counts", {}) or {}
    mode = rep.get("mode", "")
    safety = rep.get("restoration_safety", "")
    mapping = rep.get("mapping_encrypted", "")
    keyf = rep.get("key_file") or "（由 --passphrase 派生，salt 存于密钥目录）"
    skipped_detail = rep.get("skipped_detail", []) or []
    ocr_list = []
    if ocr_dir and os.path.isdir(ocr_dir):
        for root, _, files in os.walk(ocr_dir):
            for fn in files:
                ocr_list.append(os.path.join(root, fn))
    L = ["# 上云前自检报告（务必阅读并校对后再上云）", "",
         "- 生成时间：%s" % datetime.now().isoformat(timespec="seconds"),
         "- 脱敏模式：%s（恢复安全性：%s）" % (mode, safety),
         "- 命中统计：%s" % (json.dumps(counts, ensure_ascii=False) if counts else "无"),
         "", "## 一、请校对：OCR 待校对副本（未脱敏，严禁上云）",
         "位置：`02_未脱敏副本/OCR待校对/`", ""]
    if ocr_list:
        for p in ocr_list:
            L.append("- `%s`" % os.path.abspath(p))
        L.append("")
        L.append("> ⚠ 请逐字核对上述 OCR 文本，重点：姓名 / 证件号 / 银行卡 / 金额。OCR 可能误识、漏识，")
        L.append("> 确认无误后再脱敏；如发现错误，请指出以便补充 / 修正，勿直接上云。")
    else:
        L.append("- 本次无 OCR 产物，跳过。")
    L.append("")
    L.append("## 二、请复核：脱敏副本（★可上云）")
    L.append("位置：`03_脱敏副本/`")
    L.append("")
    L.append("> 请抽查脱敏是否到位（敏感字段已被替换），同时确认未误伤业务实质内容。")
    L.append("> 仅此目录可上传云端。")
    L.append("")
    L.append("## 三、本地处理·无需外发豁免留痕（如适用）")
    loc = [d for d in skipped_detail if d.get("category") == "local_only"]
    if loc:
        for d in loc:
            L.append("- `%s` — %s" % (d.get("file"), d.get("reason", "")))
        L.append("")
        L.append("> ⚠ 上述文件经你声明仅本地处理而跳过脱敏，为未脱敏原文。请确认：")
        L.append("> ① 它们已隔离存放于 `07_本地处理不外发/`（不在可上云目录）；")
        L.append("> ② 绝不外发、绝不与 `03_脱敏副本/` 一起上传。误上传会导致敏感外泄。")
    else:
        L.append("- 本次无本地处理·无需外发豁免。")
    L.append("")
    L.append("## 四、公开声明豁免留痕（如适用）")
    pub = [d for d in skipped_detail if d.get("category") == "public_declared"]
    if pub:
        for d in pub:
            L.append("- `%s` — %s" % (d.get("file"), d.get("reason", "")))
        L.append("")
        L.append("> ⚠ 上述文件经你声明为公开而跳过脱敏。请再次确认确为公开信息；误声明会导致敏感外传。")
    else:
        L.append("- 本次无公开声明豁免。")
    L.append("")
    L.append("## 五、映射表（重识别钥匙 · 保密）")
    L.append("- 加密映射表：`%s`" % mapping)
    L.append("- 密钥：`%s`（权限 600）" % keyf)
    L.append("> 映射表与副本须分离、加密、最小权限；**绝不随副本上传**。")
    L.append("")
    L.append("## 六、确认闸门（上云前必须完成）")
    L.append("请在上传云端前，确认以下各项均已完成：")
    L.append("1. ✅ 已逐字校对 OCR 待校对副本（第一节）；")
    L.append("2. ✅ 已抽查脱敏副本脱敏到位、未误伤业务（第二节）；")
    L.append("3. ✅ 已确认「本地处理·无需外发」文件隔离存放、绝不外发（第三节）；")
    L.append("4. ✅ 已阅读本报告、确认公开豁免留痕无误（第四节）。")
    L.append("")
    L.append("> **确认方式**：请明确回复「确认上云」。AI Agent 收到后才会将 `03_脱敏副本/` 上传云端；")
    L.append("> 期间如发现任何异常 / 遗漏，请直接指出，AI 将补充或修正后再请你确认。")
    L.append("> 禁止「一键脱敏即上云」。")
    with open(out, "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")
    return out

# ---------------------------------------------------------------------------
# 公开声明豁免（替代“自动识别公开主体白名单”）：默认全脱敏，用户显式声明才豁免
# ---------------------------------------------------------------------------
# 显式声明豁免——仅接受“用户明确声明”，不做文件名/目录名隐式推断
# （避免 sample/demo/pub 等常见命名、或“公开”子串被误判为公开而跳过脱敏）。
# 声明通道（见 is_public_declared，判定顺序）：
#   ① --assume-public                 提示词声明：本次输入整体公开
#   ② --public-paths <路径> [<路径>]  指定文件/文件夹公开；文件夹=其下全部（递归）
#   ③ 伴随清单 .nodesens / desensitize_manifest.json（重复使用的高级选项）
# 数据库/知识库等非文件源由 AI Agent 在闸门直接记为“已声明公开”，不经本脚本处理。
PUBLIC_MANIFEST_NAMES = (".nodesens", "desensitize_manifest.json")  # 伴随清单（高级）
PUBLIC_DECL_WARNING = ("请确认确为公开信息：被声明豁免的文件不会被脱敏，"
                       "若声明有误可能导致敏感信息未经脱敏直接外传。")
LOCAL_ONLY_WARNING = ("请确认仅本地处理、绝不外发：被声明仅本地处理的文件不会被脱敏，"
                      "若误上传云端将导致敏感信息未经脱敏直接外泄；"
                      "且不得与脱敏副本（可上云）一起存放。"
                      "本地处理中追加/生成的数据如需外发，须重新 scan 检验 + run 脱敏 + 留痕"
                      "（豁免以「不外发」为前提，外发即失效）。")


def _read_manifest(fp):
    """读取伴随清单：每行一个相对路径/通配符（fnmatch），# 开头或空行忽略。"""
    pats = []
    try:
        with open(fp, "r", encoding="utf-8") as f:
            for ln in f:
                s = ln.strip()
                if s and not s.startswith("#"):
                    pats.append(s)
    except Exception:
        pass
    return pats


def discover_public_manifests(input_path, manifest_override=None):
    """发现输入目录树内的伴随清单（.nodesens / desensitize_manifest.json），
    并合并 --public-manifest 自定义清单。
    返回 [(dir_abs, [patterns...], manifest_file), ...]；patterns 相对 dir 的路径/
    通配符，亦按 basename 匹配。自定义清单以输入根目录为基准匹配。"""
    manifests = []
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(os.path.abspath(input_path))
    src_root = os.path.abspath(src_root)
    if manifest_override and os.path.isfile(manifest_override):
        manifests.append((src_root, _read_manifest(manifest_override),
                          os.path.abspath(manifest_override)))
    if os.path.isfile(input_path):
        base = os.path.dirname(os.path.abspath(input_path))
        for name in PUBLIC_MANIFEST_NAMES:
            fp = os.path.join(base, name)
            if os.path.isfile(fp):
                manifests.append((base, _read_manifest(fp), fp))
        return manifests
    root = os.path.abspath(input_path)
    for dirpath, _dirs, filenames in os.walk(root):
        for name in PUBLIC_MANIFEST_NAMES:
            if name in filenames:
                fp = os.path.join(dirpath, name)
                manifests.append((dirpath, _read_manifest(fp), fp))
    return manifests


def _resolve_public_paths(public_paths):
    """解析 --public-paths：绝对化并分类为 [(abspath, 'dir'|'file'), ...]。
    不存在的路径忽略（偏向“有疑即脱敏”，声明错了最多多脱敏、不会漏脱敏）。"""
    out = []
    if not public_paths:
        return out
    for p in public_paths:
        ap = os.path.abspath(os.path.expanduser(p))
        if os.path.isdir(ap):
            out.append((ap, "dir"))
        elif os.path.isfile(ap):
            out.append((ap, "file"))
    return out


def is_public_declared(path, src_root, manifests, assume_public, public_paths=None):
    """返回声明豁免原因字符串；非公开返回 None。
    判定顺序：① --assume-public 全局声明；② --public-paths 指定文件/文件夹；
    ③ 伴随清单(.nodesens)匹配（高级重复使用）。
    不做任何文件名/目录名隐式推断（sample/demo/pub/“公开”等命名不再触发豁免）。"""
    if assume_public:
        return "用户声明公开(--assume-public)"
    absp = os.path.abspath(path)
    for pp, kind in _resolve_public_paths(public_paths):
        if kind == "file" and absp == pp:
            return "用户声明公开(文件: %s)" % os.path.relpath(pp, src_root)
        if kind == "dir" and (absp == pp or absp.startswith(pp + os.sep)):
            return "用户声明公开(文件夹: %s)" % os.path.relpath(pp, src_root)
    # 伴随清单（高级重复使用）
    for mdir, pats, mfile in manifests:
        mrel = os.path.relpath(absp, mdir)
        base = os.path.basename(absp)
        for pat in pats:
            if mrel == pat or fnmatch.fnmatch(mrel, pat) or fnmatch.fnmatch(base, pat):
                return "伴随清单声明(%s)" % os.path.relpath(mfile, src_root)
    return None


def is_local_declared(path, src_root, local_only, local_paths=None):
    """返回“本地处理·无需外发豁免”原因字符串；未声明返回 None。
    语义：内容可能含敏感信息，但仅本地处理、无需外发 → 跳过脱敏（不脱敏、
    不写入可上云目录、绝不外发）。与“公开声明豁免”的区别：公开=内容公开可上云；
    本地=内容保密但只本地处理，故跳过脱敏但**严禁上云**。
    判定顺序：① --local-only 全局声明；② --local-paths 指定文件/文件夹
    （文件夹=其下全部递归，文件=该文件）。不做任何文件名/目录名隐式推断。"""
    if local_only:
        return "用户声明仅本地处理(--local-only)"
    absp = os.path.abspath(path)
    for lp, kind in _resolve_public_paths(local_paths):  # 纯路径解析，public/local 共用
        if kind == "file" and absp == lp:
            return "用户声明仅本地处理(文件: %s)" % os.path.relpath(lp, src_root)
        if kind == "dir" and (absp == lp or absp.startswith(lp + os.sep)):
            return "用户声明仅本地处理(文件夹: %s)" % os.path.relpath(lp, src_root)
    return None

# 本地 OCR：rapidocr（纯 pip 安装，模型随 wheel 捆绑，完全离线、数据不出本机）。
# 引擎按需惰性初始化（首次加载模型约数秒，之后常驻复用）。
OCR_NOTE = ("本地 OCR 使用 rapidocr + onnxruntime（纯 pip、模型内置、完全离线）。"
            "若导入失败，请在本环境安装：uv add rapidocr onnxruntime pypdfium2"
            "（或 pip install rapidocr onnxruntime pypdfium2）。")

# 预处理分类 → 动作/警告 映射（供 preprocess 与 skip 复用）
PREPROCESS_ACTIONS = {
    "treatable":   ("无需预处理", None),
    "encrypted":   ("本地解密（用户提供密码，仅本地操作）", WARN_ENCRYPTED),
    "image":       ("本地 rapidocr 识别为文本并落盘，提醒用户校对", WARN_IMAGE),
    "image_pdf":   ("本地 rapidocr 识别为文本并落盘，提醒用户校对", WARN_IMAGE),
    "no_lib":      ("安装缺失的抽取库后重跑", WARN_GENERIC),
    "error":       ("人工复核：解密/确认完整性或转文本", WARN_GENERIC),
    "empty":       ("人工确认是否含敏感信息（空文档/图片型）", WARN_GENERIC),
    "unsupported": ("本地转为文本，或 OCR/解密后纳入", WARN_GENERIC),
}

# ---------------------------------------------------------------------------
# 2. 脱敏替换策略
# ---------------------------------------------------------------------------
# 生僻字 / 特殊字符判定（脱敏时优先 mask，避免保留高识别度生僻字或破坏格式的特殊字符）
# - 生僻字：CJK 扩展区（Ext-A/B/C/D/E/F/G）+ 兼容表意文字。这些字在常用字体/输入法里
#   极少出现，保留首字（如「㐀」作姓氏）会泄露高识别度特征。
# - 特殊字符：非「常见字符」——既非 ASCII 字母数字/常见符号，也非 CJK 基本区汉字、
#   常见中文标点、空白（emoji、控制字符、私有区、数学符号等）。
# - 边界：CJK 基本区（U+4E00-9FFF）内的罕见字（如「爨」「龘」）不在自动判定范围，
#   需人工复核（见 SKILL.md「生僻字与特殊字符」边界说明）。
_CJK_RARE_RANGES = (
    (0x3400, 0x4DBF),    # CJK 扩展 A（6592 生僻字）
    (0xF900, 0xFAFF),    # CJK 兼容表意文字（含罕见兼容字）
    (0x20000, 0x2FA1F),  # CJK 扩展 B~F（含兼容补充 2F800-2FA1F）
    (0x30000, 0x323AF),  # CJK 扩展 G（最新增补，极罕见）
)
_COMMON_CJK_PUNCT = set("，。、；：？！“”‘’（）【】《》—…·　")


def _is_rare_cjk(ch: str) -> bool:
    """是否 CJK 生僻字（扩展区/兼容区）。基本区罕见字（如「爨」）不在此自动判定范围。"""
    o = ord(ch)
    return any(lo <= o <= hi for lo, hi in _CJK_RARE_RANGES)


def _is_common_char(ch: str) -> bool:
    """是否「常见字符」：ASCII 字母数字/常见符号、CJK 基本区汉字、常见中文标点、空白。"""
    o = ord(ch)
    if 0x4E00 <= o <= 0x9FFF:          # CJK 基本区汉字
        return True
    if ch.isascii():
        return ch.isalnum() or ch in " -_@.:/\\"
    if ch in _COMMON_CJK_PUNCT or ch.isspace():
        return True
    return False


def _has_rare_or_special(value: str) -> bool:
    """字段值是否含生僻字（CJK 扩展/兼容区）或特殊字符（非常见字符）。
    命中 → 脱敏时优先全掩码，不保留含此类字符的首字/尾字。"""
    return any(_is_rare_cjk(c) or not _is_common_char(c) for c in value)


# 「前 N 位 + 星号 + 后 M 位」同构掩码规则表：kind -> (长度阈值, 保留前, 保留后, 是否精确等长)
# 精确等长（True）用于正则本身定长的字段（手机 11 / 身份证 18）；其余用最短长度（>=）。
_KEEP_MASK_RULES = {
    "phone":     (11, 3, 4, True),
    "id_card":   (18, 6, 4, True),
    "bank_card": (12, 0, 4, False),
    "plate":     (7,  2, 2, False),
    "passport":  (8,  1, 2, False),
    "iban":      (15, 4, 4, False),
    "swift":     (8,  4, 2, False),
    "vat":       (10, 2, 4, False),
}


def _keep_mask(value, keep_l, keep_r):
    """保留前 keep_l 位与后 keep_r 位，中间以星号掩码。"""
    return value[:keep_l] + "*" * (len(value) - keep_l - keep_r) + value[-keep_r:]


def mask_value(kind: str, value: str) -> str:
    """按类型做掩码（可逆映射由调用方记录）。"""
    rule = _KEEP_MASK_RULES.get(kind)
    if rule is not None:
        min_len, keep_l, keep_r, exact = rule
        if (len(value) == min_len if exact else len(value) >= min_len):
            return _keep_mask(value, keep_l, keep_r)
    if kind == "email":
        local, _, domain = value.partition("@")
        keep = local[0] if local else "x"
        return keep + "***@" + domain
    if kind == "ip":
        parts = value.split(".")
        if len(parts) == 4:
            return parts[0] + "." + parts[1] + ".*.*"
        return "*.**.**.*"
    if kind == "jwt":
        # 保留 eyJ 前缀与末 6 位（便于人工辨认与恢复），中间掩码
        if len(value) > 12:
            return value[:6] + "***" + value[-6:]
        return value[:3] + "***"
    if kind == "intl_phone":
        # 保留 + 号与末 4 位数字（国家码长度 1~3 位不定，不硬拆，避免误导）
        digits = re.sub(r"\D", "", value)
        if len(digits) >= 6:
            return "+" + "*" * (len(digits) - 4) + digits[-4:]
        return "+" + "*" * max(len(digits) - 2, 1)
    if kind in ("cn_name", "name"):
        # 含生僻字/特殊字符 → 优先全掩码（避免保留高识别度生僻字或破坏格式的特殊字符）
        if _has_rare_or_special(value):
            return "*" * len(value)
        # 英文姓名（含空格，如 John Smith）→ 按词保留首字母（J*** S****）
        if " " in value:
            return " ".join(w[0] + "*" * (len(w) - 1) for w in value.split() if w)
        # 中文姓名 → 保留姓氏/首字，其余脱敏（原始值仍记录于加密映射表，可逆）
        return value[0] + "*" * (len(value) - 1)
    if kind == "cn_address":
        return "[地址]"
    if kind == "cn_org":
        return "[机构]"
    # 兜底：含生僻字/特殊字符 → 全掩码；否则首尾各留 1，中间掩码
    if _has_rare_or_special(value):
        return "*" * len(value)
    if len(value) <= 2:
        return "*" * len(value)
    return value[0] + "*" * (len(value) - 2) + value[-1]


def redact_value(kind: str, value: str) -> str:
    return "[REDACTED_%s]" % kind.upper()


def token_value(token_map: dict, value: str) -> str:
    """为同一原始值生成稳定令牌（跨记录可关联）。"""
    if value in token_map:
        return token_map[value]
    tok = "T%04d" % (len(token_map) + 1)
    token_map[value] = tok
    return tok


# hybrid 模式分隔符：语义掩码 + 唯一令牌，形如  张*⟦T0001⟧
# 选用数学空心方括号（U+27E6 / U+27E7），在正常数据中极罕见，避免与正文混淆。
HYBRID_SEP_L, HYBRID_SEP_R = "⟦", "⟧"


# ---------------------------------------------------------------------------
# 3. 识别与脱敏核心
# ---------------------------------------------------------------------------
def desensitize_text(text: str, mode: str, token_map: dict, counts: dict,
                     names: set = None, patterns=None, custom_map=None):
    """对单段文本做脱敏，返回 (脱敏后文本, 命中列表)。

    注意：必须一次性构建新串（re.sub + 回调），不可在 finditer 迭代中
    原地切片替换——否则变长替换会导致后续匹配偏移错位。
    替换串长度完全可变（可长可短，甚至远长于原文，如 hybrid 的“掩码⟦Txxxx⟧”）：
    re.sub 会按需动态拼接新串，无需、也无法在扫描阶段预分配固定长度，
    因此不存在“替换串长度不足”的失败模式——令牌超过 9999 个时自动扩展为
    T10000… 仍唯一且可被分隔符 ⟦ ⟧ 正确解析。

    custom_map：用户自定义脱敏映射（原始值→替换值），**独立于识别正则**——
    用户指定即视为需脱敏，精确匹配后替换为用户指定值（原样、不加令牌），
    可逆性由加密映射表反查保证。
    """
    if patterns is None:
        patterns = PATTERNS
    hits = []
    # 全角数字/字母归一为半角（全角手机/身份证/银行卡/IBAN 等才能被正则识别）
    text = _normalize_fullwidth(text)
    result = text

    # 3.0 用户自定义脱敏映射（原始值→替换值）：主动精确匹配，命中即替换为用户指定值
    #     （先于 names/正则，故用户映射的值不会再次被自动掩码覆盖）
    if custom_map:
        for orig in sorted(custom_map, key=len, reverse=True):
            if not orig:
                continue
            repl = custom_map[orig]
            def _cb(m, _orig=orig, _repl=repl):
                counts["custom"] = counts.get("custom", 0) + 1
                hits.append({"type": "custom", "original": _orig, "replacement": _repl})
                return _repl
            result = re.sub(re.escape(orig), _cb, result)

    # 3.1 自定义姓名清单（若有）
    if names:
        name_cb = _repl_closure("name", mode, token_map, counts, hits)
        for nm in sorted(names, key=len, reverse=True):
            if not nm:
                continue
            result = re.sub(re.escape(nm), name_cb, result)

    # 3.2 顺序应用正则（先处理 ID/手机，避免银行卡重复命中；跨境标识在数字类之后，互不冲突）
    for kind in ["id_card", "phone", "bank_card", "ip", "email", "jwt",
                 "plate", "passport", "iban", "swift", "vat", "intl_phone"]:
        pat = patterns[kind]
        cb = _repl_closure(kind, mode, token_map, counts, hits)
        result = pat.sub(lambda m, _cb=cb: _cb(m), result)

    # 3.3 中文增强（cn_name 为双模式列表，其余为单正则）
    if "cn_name" in patterns:
        for npat in patterns["cn_name"]:
            cb = _repl_closure("cn_name", mode, token_map, counts, hits)
            result = npat.sub(lambda m, _cb=cb: _cb(m), result)
    if "cn_address" in patterns:
        cb = _repl_closure("cn_address", mode, token_map, counts, hits)
        result = patterns["cn_address"].sub(lambda m, _cb=cb: _cb(m), result)
    if "cn_org" in patterns:
        cb = _repl_closure("cn_org", mode, token_map, counts, hits)
        result = patterns["cn_org"].sub(lambda m, _cb=cb: _cb(m), result)

    return result, hits


def _repl_closure(kind: str, mode: str, token_map: dict, counts: dict, hits: list):
    """生成脱敏回调：记录命中并返回替换串。"""
    def cb(m):
        seg = m.group(0)
        # 已被 ID/手机占位覆盖的纯数字段不应再当银行卡
        if kind == "bank_card" and set(seg) <= set("*"):
            return seg
        # 跨境标识（iban/swift/vat）须经校验（MOD-97 / ISO 国家码），校验失败不脱敏（降误报）
        validator = _KIND_VALIDATOR.get(kind)
        if validator is not None and not validator(seg):
            return seg
        if mode == "redact":
            repl = redact_value(kind, seg)
        elif mode == "token":
            repl = token_value(token_map, seg)
        elif mode == "hybrid":
            # 语义掩码 + 唯一令牌：保留字段类型，同时靠令牌实现无歧义恢复
            repl = "%s%s%s%s" % (mask_value(kind, seg), HYBRID_SEP_L,
                                 token_value(token_map, seg), HYBRID_SEP_R)
        else:
            repl = mask_value(kind, seg)
        counts[kind] = counts.get(kind, 0) + 1
        hits.append({"type": kind, "original": seg, "replacement": repl})
        return repl
    return cb


def scan_text(text: str, names: set = None, patterns=None):
    """仅扫描，返回命中计数（复用脱敏的顺序逻辑，保证与 run 一致）。"""
    token_map = {}
    counts = {}
    desensitize_text(text, "mask", token_map, counts, names, patterns)
    return counts


def _read_text(path):
    """读取文本文件并自动探测编码：BOM → UTF-8 → charset_normalizer → 兜底。

    电商/小微/跨境导出文件编码混杂：国内报表常为 GBK/GB2312，日本乐天/Amazon JP 为
    Shift-JIS(cp932)，港台繁体为 BIG5，Windows 拉丁为 cp1252。此前仅 UTF-8→GB18030，
    日文/繁体/拉丁重音会乱码（非 ASCII PII 不可读、无法脱敏）。现用 charset_normalizer
    （纯 Python、离线，作为 requests 的传递依赖已随装）探测编码族，覆盖 cp932/big5/
    cp1252/gb18030 等；探测不可用时回退 UTF-8→GB18030→cp1252→replace。"""
    with open(path, "rb") as f:
        raw = f.read()
    # 1) BOM 优先（UTF-8-SIG 等带 BOM 文件直接按标记解码）
    if raw.startswith(b"\xef\xbb\xbf"):
        return raw.decode("utf-8-sig")
    if raw.startswith(b"\xff\xfe"):
        return raw.decode("utf-16-le")
    if raw.startswith(b"\xfe\xff"):
        return raw.decode("utf-16-be")
    # 2) UTF-8 严格解码（现代默认，误判率极低）
    try:
        return raw.decode("utf-8")
    except (UnicodeDecodeError, LookupError):
        pass
    # 3) charset_normalizer 探测编码族（cp932/big5/cp1252/gb18030 等）
    try:
        from charset_normalizer import from_bytes
        best = from_bytes(raw).best()
        if best is not None and best.encoding:
            try:
                return raw.decode(best.encoding)
            except (UnicodeDecodeError, LookupError):
                pass
    except Exception:
        pass
    # 4) 兜底：GB18030 → cp1252（cp1252 几乎总能解码，仅作最后兜底）
    for enc in ("gb18030", "cp1252"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def _write_text_raw(path: str, text: str) -> None:
    """写文本文件（newline=""：保留原始换行，Windows 不把 \\n 二次翻译为 \\r\\n）。

    与 _read_text（二进制读、保留 \\r\\n）对称，保证「读保留、写不翻译」，
    使脱敏副本/还原文件与原文件逐字节一致。修复 Windows 下 CRLF 双倍换行 bug
    （此前写出走默认文本模式，Windows 上 \\r\\n 被翻成 \\r\\n\\r\\n）。
    """
    with open(path, "w", encoding="utf-8", newline="") as f:
        f.write(text)


def _read_text_plain(path: str) -> str:
    """读文本文件（newline=""：不做 universal-newline 翻译，保留原始 \\r\\n）。

    供 restore 读取脱敏副本，避免 Windows 读时把 \\r\\n 折叠为 \\n 后再写出造成二次翻译。
    """
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        return f.read()


def process_file(path: str, mode: str, out_root: str, src_root: str, token_map: dict,
                 mapping: dict, names: set, counts_total: dict, do_write: bool,
                 patterns=None, skipped=None, public_manifests=None,
                 assume_public=False, public_paths=None, local_only=False,
                 local_paths=None, local_out_root=None, cleaning=None,
                 custom_map=None, tabular_names=False):
    """处理单个文件。skipped 为可选 list，用于收集“未处理文件”（skip manifest）。
    public_manifests/assume_public/public_paths 用于“公开声明豁免”：命中则跳过脱敏
    （留痕+警示，绝不静默）。local_only/local_paths 用于“本地处理·无需外发豁免”：
    命中则跳过脱敏且绝不写入可上云目录（工作区模式下原样复制到隔离目录 07_本地处理不外发/，
    与 03_脱敏副本/ 物理隔离）。豁免不做文件名/目录名隐式推断。
    二者同时命中时「本地处理」优先（不外发比可上云更严格）。
    custom_map 为用户自定义脱敏映射（原始值→替换值），命中优先用用户替换值。
    cleaning 为可选 dict，累计“疑似未清洗数据形态”计数（仅提醒，不脱敏）。"""
    try:
        # 跳过流程自身产出的元数据文件，避免把清单/报告当作业务文本脱敏
        if os.path.basename(path) in META_SKIP:
            return "skip_meta"
        # 豁免判定（本地处理优先于公开：不外发比可上云更严格，先判 local）
        local_reason = is_local_declared(path, src_root, local_only, local_paths)
        pub_reason = is_public_declared(path, src_root, public_manifests,
                                        assume_public, public_paths)
        if local_reason:
            if do_write:
                # run 模式：仅本地处理不脱敏、绝不写入可上云目录；工作区模式复制到隔离目录留痕
                if local_out_root:
                    _copy_to_local_only(path, src_root, local_out_root)
                _record_local_only(skipped, path, src_root,
                                   local_reason + "，已跳过脱敏（仅本地处理，严禁外发）")
                return "skipped"
            # scan 模式：继续下方读取并检测（fail-safe：仍报命中，不计入脱敏计划）
        elif pub_reason:
            if do_write:
                # run 模式：被声明公开的文件不脱敏、不写入输出目录（仅留痕）
                _record_public_declared(skipped, path, src_root,
                                        pub_reason + "，已跳过脱敏")
                return "skipped"
            # scan 模式：继续下方读取并检测（fail-safe：仍报命中，不计入脱敏计划）
        ext = os.path.splitext(path)[1].lower()
        # 二进制/库文件：抽取文本后再脱敏（库缺失则跳过并告警）
        if ext in EXTRACT_EXTS:
            return _process_extracted(path, mode, out_root, src_root, token_map,
                                      mapping, names, counts_total, do_write,
                                      patterns, skipped, local_reason, pub_reason,
                                      local_out_root, custom_map,
                                      tabular_names=tabular_names)
        # 既非文本/代码也非可抽取类型：登记为未处理（不静默忽略）
        if ext not in (TEXT_EXTS | CODE_EXTS):
            if skipped is not None:
                if local_reason:
                    # 声明仅本地处理但非文本类：登记为 local_only（覆盖默认分类）
                    if local_out_root:
                        _copy_to_local_only(path, src_root, local_out_root)
                    _record_local_only(skipped, path, src_root,
                                       local_reason + "，已跳过脱敏（非文本类，未检测，仅本地处理）")
                elif pub_reason:
                    # 声明公开但非文本类：登记为公开声明（覆盖默认分类）
                    _record_public_declared(skipped, path, src_root,
                                            pub_reason + "，已跳过脱敏（非文本类，未检测）")
                elif ext in IMAGE_EXTS:
                    _skip(skipped, path, src_root, "image",
                          "影像/图片文件（%s）：scan/run 不直接做 OCR，请先运行 preprocess"
                          "（内置 rapidocr 本地 OCR，完全离线）识别为文本后再纳入脱敏；"
                          "原始图片不会被处理" % ext)
                else:
                    _skip(skipped, path, src_root, "unsupported",
                          "未支持的信息源类型（%s）：二进制/未知格式，若含敏感文本请先转为"
                          "文本文件（或 OCR/解密）后再纳入脱敏" % ext)
            return "skipped"
        text = _read_text(path)
        # 结构化列解析（--tabular-names）：csv/tsv 解析列头，识别姓名列，收集英文姓名并入 names
        if tabular_names and ext in (".csv", ".tsv"):
            _tab = _collect_tabular_names(text, "," if ext == ".csv" else "\t")
            if _tab:
                names = (names or set()) | _tab
    except Exception as e:
        print("  [跳过] 无法读取 %s: %s" % (path, e), file=sys.stderr)
        _skip(skipped, path, src_root, "error", "读取失败: %s" % e)
        return "error"

    ext = os.path.splitext(path)[1].lower()
    # 密钥检测扩展：html/htm 一并启用（内嵌 <script> 的 adminToken/secret 赋值是真实泄漏点）
    is_code = ext in SECRET_EXTS
    rel_path = os.path.relpath(path, src_root)

    # 清洗建议：检测疑似「未清洗/不规范」数据形态（仅提醒、不脱敏，见 CLEANING_ADVICE_TEXT）
    adv = _find_cleaning_advice(text)
    if adv and cleaning is not None:
        for k, v in adv.items():
            cleaning[k] = cleaning.get(k, 0) + v

    if not do_write:
        c = scan_text(text, names, patterns)
        if is_code:
            sec = len(SECRET_PATTERN.findall(text))
            if sec:
                c["secret"] = c.get("secret", 0) + sec
            # SQL INSERT 位置参数口令（与 run 分支保持一致计数）
            if ext == ".sql":
                sec2 = len(_sql_insert_secret_values(text))
                if sec2:
                    c["secret"] = c.get("secret", 0) + sec2
        if local_reason:
            # scan 仍检测：若检出疑似敏感，追加警示；声明文件不计入脱敏计划
            reason = local_reason + "，已跳过脱敏（仅本地处理，严禁外发）"
            if c:
                reason += "（scan 仍检出 %d 处疑似敏感，请复核）" % sum(c.values())
                print("  %-40s %s  ⚠ 声明仅本地处理但检出疑似敏感" % (rel_path, c))
            _record_local_only(skipped, path, src_root, reason)
            return "skipped"
        if pub_reason:
            # scan 仍检测：若检出疑似敏感，追加警示；声明文件不计入脱敏计划
            reason = pub_reason + "，已跳过脱敏"
            if c:
                reason += "（scan 仍检出 %d 处疑似敏感，请复核）" % sum(c.values())
                print("  %-40s %s  ⚠ 声明公开但检出疑似敏感" % (rel_path, c))
            _record_public_declared(skipped, path, src_root, reason)
            return "skipped"
        if c:
            print("  %-40s %s" % (rel_path, c))
            for k, v in c.items():
                counts_total[k] = counts_total.get(k, 0) + v
        if adv:
            print("  %-40s ⚠清洗建议 %s（先清洗再脱敏）" % (rel_path, adv))
        return "processed"

    # 代码/类代码文件额外处理密钥
    if is_code:
        def _sec_repl(m):
            val = m.group("val")
            # 剥离包裹引号，记录纯净原始值
            if len(val) >= 2 and val[0] in "'\"" and val[-1] == val[0]:
                inner = val[1:-1]
            else:
                inner = val
            counts_total["secret"] = counts_total.get("secret", 0) + 1
            if mode == "hybrid":
                repl = "***%s%s%s" % (HYBRID_SEP_L, token_value(token_map, inner),
                                      HYBRID_SEP_R)
            else:
                repl = "***"
            mapping.setdefault(rel_path, []).append(
                {"type": "secret", "original": inner, "replacement": repl})
            return m.group(0).replace(val, repl)
        text = SECRET_PATTERN.sub(_sec_repl, text)
        # SQL INSERT 位置参数口令：VALUES(..., 'p@ssw0rd', ...) 前无 password= 前缀，
        # SECRET_PATTERN 捕获不到，按"列名含口令关键词 → 对应位置值"定位脱敏。
        if ext == ".sql":
            for quoted, inner in _sql_insert_secret_values(text):
                counts_total["secret"] = counts_total.get("secret", 0) + 1
                if mode == "hybrid":
                    repl = "***%s%s%s" % (HYBRID_SEP_L, token_value(token_map, inner),
                                          HYBRID_SEP_R)
                else:
                    repl = "***"
                mapping.setdefault(rel_path, []).append(
                    {"type": "secret", "original": inner, "replacement": repl})
                text = text.replace(quoted, repl, 1)

    new_text, hits = desensitize_text(text, mode, token_map, counts_total, names, patterns, custom_map)
    for h in hits:
        h["file"] = rel_path
        mapping.setdefault(rel_path, []).append(h)

    out_path = os.path.join(out_root, rel_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    _write_text_raw(out_path, new_text)
    return "processed"


# ---------------------------------------------------------------------------
# 3.1 二进制/库文件文本抽取（抽取后再脱敏；库缺失则跳过并告警）
# ---------------------------------------------------------------------------
def _lib_for(ext):
    return {".docx": "python-docx", ".xlsx": "openpyxl",
            ".pptx": "python-pptx", ".pdf": "pypdfium2"}.get(ext, "未知库")


def _extract_sqlite(path):
    """用内置 sqlite3 把库表内容导出为文本（零依赖）。"""
    import sqlite3
    con = sqlite3.connect(path)
    try:
        cur = con.cursor()
        lines = []
        for (name,) in cur.execute(
                "SELECT name FROM sqlite_master WHERE type='table'"):
            try:
                rows = con.execute('SELECT * FROM "%s"' % name)
                cols = [d[0] for d in rows.description]
                lines.append("TABLE %s (%s)" % (name, ", ".join(cols)))
                for row in rows.fetchall():
                    lines.append(" | ".join("" if v is None else str(v) for v in row))
            except Exception:
                pass
        return "\n".join(lines)
    finally:
        con.close()


# xlsx 批注二次打开的体积阈值：超过则跳过批注扫描（read_only 正文不受影响），
# 避免超大表非只读整表载入的内存开销；批注内敏感信息改由人工检查提示兜底。
_XLSX_COMMENT_LIMIT = 15 * 1024 * 1024  # 15MB


def _extract_text(ext, path):
    """返回 (文本, 缺失库名或None, 原因或None)。

    text=None 表示未成功抽取。原因 reason 取值：
      None      成功抽取到文本
      'no_lib'  缺少第三方抽取库（ImportError）
      'encrypted' 文件加密（异常信息含 password/encrypt/decrypt/protected）
      'image_only' PDF 抽到空文本（纯图片型/扫描件 PDF）
      'empty'    其他文档抽到空文本
      'error'    其他抽取异常（可能加密/损坏）
    """
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            # 段落 + 表格都要抽：小微常用 Word 表格做登记表（发票/花名册/台账），
            # 此前只抽段落导致表格内身份证/税号/金额漏检。
            parts = [p.text for p in doc.paragraphs if p.text]
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text for c in row.cells if c.text]
                    if cells:
                        parts.append(" | ".join(cells))
            text = "\n".join(parts)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            # 正文：read_only 快速路径（超大表内存友好、不整表载入）。
            wb = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        lines.append(" ".join(cells))
            try:
                wb.close()
            except Exception:
                pass
            # 批注：非只读**二次打开**单独扫描（read_only 模式读不到批注）。
            # 大文件非只读加载内存开销高：超过阈值则跳过批注扫描并明确提示
            # （批注多为少量补充信息，如"客户手机 … 记得回访"，提示后由用户人工检查）。
            if os.path.getsize(path) <= _XLSX_COMMENT_LIMIT:
                wb2 = load_workbook(path, read_only=False, data_only=True)
                try:
                    for ws in wb2.worksheets:
                        for row in ws.iter_rows():
                            for cell in row:
                                if cell.comment and cell.comment.text:
                                    lines.append("[批注] %s" % cell.comment.text)
                finally:
                    try:
                        wb2.close()
                    except Exception:
                        pass
            else:
                print("  [提示] %s 体积 >%dMB，跳过批注扫描：批注内如有敏感信息（手机/证件等）"
                      "请人工检查后再上云。" % (os.path.basename(path),
                                                _XLSX_COMMENT_LIMIT // (1024 * 1024)),
                      file=sys.stderr)
            text = "\n".join(lines)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(r.text for r in para.runs)
                            if t:
                                lines.append(t)
            text = "\n".join(lines)
        elif ext == ".pdf":
            # pypdfium2（Chrome 同款 PDFium，C 实现）：快且对畸形 PDF 更稳健；
            # 加密 PDF 抛 PdfiumError → 由下方异常分支归类 encrypted
            import pypdfium2 as pdfium
            doc = pdfium.PdfDocument(path)
            try:
                text = "\n".join(
                    doc[i].get_textpage().get_text_range()
                    for i in range(len(doc)))
            finally:
                doc.close()
        elif ext == ".db":
            text = _extract_sqlite(path)
        else:
            return (None, None, "error")
    except ImportError:
        return (None, _lib_for(ext), "no_lib")
    except Exception as e:
        msg = str(e).lower()
        if any(k in msg for k in ("password", "encrypt", "decrypt", "protected",
                                  "not a database")):
            return (None, None, "encrypted")
        # pypdfium2 对真加密 PDF 抛 PdfiumError（消息可能不含关键字）；
        # 此时用 /Encrypt 标记二次确认，避免与“文件损坏”混淆
        if ext == ".pdf" and _pdf_has_encrypt(path):
            return (None, None, "encrypted")
        return (None, None, "error")
    # 抽到文本但为空：可能是纯图片型扫描件，也可能是加密 PDF
    # （个别加密 PDF 可开但无文本层，故额外探测 /Encrypt 标记）
    if not text or not text.strip():
        if ext == ".pdf" and _pdf_has_encrypt(path):
            return (None, None, "encrypted")
        return (None, None, "image_only" if ext == ".pdf" else "empty")
    return (text, None, None)


def _pdf_has_encrypt(path):
    """探测 PDF 是否含 /Encrypt 字典（粗略但高效）。

    /Encrypt 引用通常位于文件**尾部**的 trailer，但也可能出现在前部对象中，
    故头、尾各读一段（各至多 5MB），避免大文件只读头部造成漏判。"""
    try:
        size = os.path.getsize(path)
        chunk = 5 * 1024 * 1024
        with open(path, "rb") as fb:
            head = fb.read(chunk)
            if b"/Encrypt" in head:
                return True
            if size > chunk:
                fb.seek(max(0, size - chunk))
                return b"/Encrypt" in fb.read()
        return False
    except Exception:
        return False


def _process_extracted(path, mode, out_root, src_root, token_map, mapping,
                       names, counts_total, do_write, patterns, skipped,
                       local_reason=None, pub_reason=None, local_out_root=None,
                       cleaning=None, custom_map=None, tabular_names=False):
    ext = os.path.splitext(path)[1].lower()
    rel_path = os.path.relpath(path, src_root)
    text, missing, reason = _extract_text(ext, path)
    # 结构化列解析（--tabular-names）：xlsx 用 openpyxl 读列头收集英文姓名（复用 names 机制）
    if tabular_names and ext == ".xlsx" and names is not None:
        _tab = _collect_tabular_names_xlsx(path)
        if _tab:
            names = (names or set()) | _tab
    if text is None:
        if missing:
            cat, reason = "no_lib", ("缺少抽取库 %s，无法处理该二进制文档：请安装对应库"
                   "（docx→python-docx，xlsx→openpyxl，pptx→python-pptx，"
                   "pdf→pypdfium2）后再处理" % missing)
        else:
            cat = {"encrypted": "encrypted", "image_only": "image_pdf",
                   "empty": "empty", "error": "error"}.get(reason, "error")
            reason = REMIND.get(reason, "文本抽取失败（请确认文件未加密且未损坏）")
        if local_reason:
            # 声明仅本地处理但无法抽取文本：登记为 local_only，不按异常报
            _record_local_only(skipped, path, src_root,
                               local_reason + "，已跳过脱敏（无法抽取文本，未检测，仅本地处理）")
        elif pub_reason:
            # 声明公开但无法抽取文本（加密/图片型等）：登记为公开声明，不按异常报
            _record_public_declared(skipped, path, src_root,
                                    pub_reason + "，已跳过脱敏（无法抽取文本，未检测）")
        else:
            _skip(skipped, path, src_root, cat, reason)
            print("  [跳过] %s：%s" % (rel_path, reason), file=sys.stderr)
        return "skipped"
    # 清洗建议：二进制文档（表格/批注）同样检测疑似未清洗形态（仅提醒、不脱敏）
    adv = _find_cleaning_advice(text)
    if adv and cleaning is not None:
        for k, v in adv.items():
            cleaning[k] = cleaning.get(k, 0) + v
    if not do_write:
        c = scan_text(text, names, patterns)
        if local_reason:
            # scan 仍检测：若检出疑似敏感，追加警示；不计入脱敏计划
            reason = local_reason + "，已跳过脱敏（仅本地处理，严禁外发）"
            if c:
                reason += "（scan 仍检出 %d 处疑似敏感，请复核）" % sum(c.values())
                print("  %-40s %s  ⚠ 声明仅本地处理但检出疑似敏感" % (rel_path, c))
            _record_local_only(skipped, path, src_root, reason)
            return "skipped"
        if pub_reason:
            # scan 仍检测：若检出疑似敏感，追加警示；不计入脱敏计划
            reason = pub_reason + "，已跳过脱敏"
            if c:
                reason += "（scan 仍检出 %d 处疑似敏感，请复核）" % sum(c.values())
                print("  %-40s %s  ⚠ 声明公开但检出疑似敏感" % (rel_path, c))
            _record_public_declared(skipped, path, src_root, reason)
            return "skipped"
        if c:
            print("  %-40s %s" % (rel_path, c))
            for k, v in c.items():
                counts_total[k] = counts_total.get(k, 0) + v
        if adv:
            print("  %-40s ⚠清洗建议 %s（先清洗再脱敏）" % (rel_path, adv))
        return "processed"
    if local_reason:
        # run 模式：仅本地处理不脱敏、绝不写入可上云目录；工作区模式复制到隔离目录留痕
        if local_out_root:
            _copy_to_local_only(path, src_root, local_out_root)
        _record_local_only(skipped, path, src_root,
                           local_reason + "，已跳过脱敏（仅本地处理，严禁外发）")
        return "skipped"
    if pub_reason:
        _record_public_declared(skipped, path, src_root, pub_reason + "，已跳过脱敏")
        return "skipped"
    new_text, hits = desensitize_text(text, mode, token_map, counts_total, names, patterns, custom_map)
    for h in hits:
        h["file"] = rel_path
        mapping.setdefault(rel_path, []).append(h)
    out_path = os.path.join(out_root, rel_path) + ".txt"
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    _write_text_raw(out_path, new_text)
    return "processed"


# ---------------------------------------------------------------------------
# 4. 映射表加密
# ---------------------------------------------------------------------------
def derive_key_from_passphrase(passphrase: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=200_000)
    return urlsafe_b64(kdf.derive(passphrase.encode()))


def urlsafe_b64(raw: bytes) -> bytes:
    import base64
    return base64.urlsafe_b64encode(raw)


def make_fernet(key: bytes) -> Fernet:
    return Fernet(key)


def encrypt_mapping(mapping: dict, keys_dir: str, passphrase: str = None) -> dict:
    """加密映射表。返回 {mapping_file, key_file|None, salt_file?}。"""
    os.makedirs(keys_dir, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    fernet = None
    key_file = None
    if passphrase:
        salt = secrets.token_bytes(16)
        key = derive_key_from_passphrase(passphrase, salt)
        fernet = make_fernet(key)
        salt_path = os.path.join(keys_dir, "salt_%s.bin" % ts)
        with open(salt_path, "wb") as f:
            f.write(salt)
    else:
        key = Fernet.generate_key()
        fernet = make_fernet(key)
        key_file = os.path.join(keys_dir, "desensitize_key_%s.key" % ts)
        with open(key_file, "wb") as f:
            f.write(key)
        os.chmod(key_file, 0o600)

    payload = json.dumps(mapping, ensure_ascii=False, indent=2).encode("utf-8")
    enc = fernet.encrypt(payload)
    enc_path = os.path.join(keys_dir, "mapping_%s.json.enc" % ts)
    with open(enc_path, "wb") as f:
        f.write(enc)
    os.chmod(enc_path, 0o600)
    return {"mapping_file": enc_path, "key_file": key_file, "encrypted": True}


def find_collisions(mapping: dict):
    """唯一性检测：找出"同一脱敏值 → 多个不同原始值"的多对一歧义。

    返回 {replacement: [去重后的原始值列表]}（仅含歧义项）。
    恢复（decrypt）时，若某 replacement 对应多个 original，则无法唯一还原，存在混淆风险。
    - token 模式：每个原始值分配唯一令牌，天然不会进入此表。
    - hybrid 模式：语义掩码 + 唯一令牌，replacement 含唯一令牌，天然不会进入此表。
    - mask 模式：中文姓名/地址/机构、邮箱等同型值易碰撞，需据此告警。
    - redact 模式：不可逆，恢复本就不适用。
    """
    rev = {}
    for items in mapping.values():
        for it in items:
            rev.setdefault(it["replacement"], set()).add(it["original"])
    return {r: sorted(v) for r, v in rev.items() if len(v) > 1}


def _print_exempt_block(title, desc, items, tip, src_root):
    """打印一类豁免项（local_only / public_declared）的清单块。"""
    print("\n[%s] 以下 %d 个文件%s，已跳过脱敏（不脱敏、不计数、不报告命中）："
          % (title, len(items), desc), file=sys.stderr)
    for it in items:
        rel = it.get("rel") or os.path.relpath(it.get("path", ""), src_root)
        reason = it.get("reason", "")
        print("    - %s  （%s）" % (rel, reason), file=sys.stderr)
        sz = it.get("size")
        sh = it.get("sha256")
        if sz is not None:
            print("        🔒 %d bytes  sha256:%s" % (sz, sh), file=sys.stderr)
        w = it.get("warning")
        if w:
            print("        ⚠ %s" % w, file=sys.stderr)
    print("  → %s" % tip, file=sys.stderr)


def _report_skipped(skipped, src_root):
    """打印“未处理文件清单”（skip manifest），杜绝目录扫描的静默漏扫。
    区分三类：local_only（仅本地处理·无需外发，已跳过脱敏）、public_declared（用户显式
    声明公开，已跳过脱敏）与其他（可能含敏感，须处理）。
    兼容旧式 (path, reason) 元组与新式 dict 两种结构。"""
    if not skipped:
        return
    loc = []
    pub = []
    others = []
    for it in skipped:
        cat = it.get("category") if isinstance(it, dict) else None
        (loc if cat == "local_only"
         else pub if cat == "public_declared"
         else others).append(it)

    if loc:
        _print_exempt_block(
            "本地处理·无需外发豁免",
            "经用户声明仅本地处理（--local-only / --local-paths）",
            loc,
            "请确认这些文件仅本地处理、绝不外发；误上传云端会导致敏感信息未经脱敏直接外泄。",
            src_root)

    if pub:
        _print_exempt_block(
            "公开声明豁免",
            "经用户显式声明为公开（--assume-public / --public-paths / 伴随清单）",
            pub,
            "请确认这些文件确为公开信息；误声明会导致敏感信息未经脱敏直接外传。",
            src_root)

    if others:
        print("\n[警告] 以下 %d 个文件未被处理（可能含敏感信息，上云前须先转文本/OCR/"
              "安装抽取库，或人工复核）：" % len(others), file=sys.stderr)
        for it in others:
            if isinstance(it, dict):
                rel = it.get("rel") or os.path.relpath(it.get("path", ""), src_root)
                reason = it.get("reason", "")
                warn = it.get("warning")
                line = "    - %s  （%s）" % (rel, reason)
                print(line, file=sys.stderr)
                if warn:
                    print("        %s" % warn, file=sys.stderr)
            else:
                p, reason = it
                print("    - %s  （%s）" % (os.path.relpath(p, src_root), reason),
                      file=sys.stderr)


def report_collisions(collisions: dict, mode: str):
    """打印碰撞告警（到 stderr），并给出模式建议。"""
    if not collisions:
        return
    print("\n[警告] 检测到 %d 个脱敏值存在“多对一”歧义（恢复时可能混淆）："
          % len(collisions), file=sys.stderr)
    for r, origs in collisions.items():
        preview = " / ".join(origs[:8]) + (" …" if len(origs) > 8 else "")
        print("    %s  <-  %s" % (r, preview), file=sys.stderr)
    if mode == "mask":
        print("  → 原因：mask 模式按类型生成固定掩码，同型不同值会碰撞（如不同姓名同为"
              "\"张*\"）。", file=sys.stderr)
        print("  → 建议：如需无歧义恢复，请改用 --mode hybrid（语义掩码+唯一令牌，"
              "保留字段语义且天然 1:1）或 --mode token。", file=sys.stderr)


# ---------------------------------------------------------------------------
# 5. 命令行
# ---------------------------------------------------------------------------
def iter_targets(input_path: str, recursive: bool):
    """产出待处理文件。**注意：这里产出目录下全部文件**，由 process_file 按扩展名
    决定处理或登记为“未处理”（skip manifest），从而杜绝“静默漏扫”。"""
    if os.path.isfile(input_path):
        yield input_path
        return
    if recursive:
        for root, dirs, files in os.walk(input_path):
            dirs.sort()  # 排序保证跨平台遍历顺序确定（token 编号可复现）
            for fn in sorted(files):
                yield os.path.join(root, fn)
    else:
        for fn in sorted(os.listdir(input_path)):
            p = os.path.join(input_path, fn)
            if os.path.isfile(p):
                yield p


def load_names(path: str):
    if not path:
        return set()
    try:
        with open(path, "r", encoding="utf-8") as f:
            return {line.strip() for line in f if line.strip()}
    except FileNotFoundError:
        print("  [警告] 姓名清单文件不存在：%s（已忽略）" % path, file=sys.stderr)
        return set()


def load_custom_mapping(path: str):
    """加载用户自定义脱敏映射（原始值 → 替换值），供 `--mapping` 使用。
    支持两种格式：
      - JSON：{"原始值": "替换值", ...}
      - 文本：每行 "原始值=替换值" 或 "原始值<TAB>替换值"（# 注释、空行忽略）。
    ⚠ 该映射文件含敏感信息（原始值），严禁外传；仅本地使用，命中结果随加密
    映射表（.desensitize_keys/）本地保存、绝不外发。文件不存在/解析失败返回 None
    （调用方告警后忽略，不影响主流程）。"""
    if not path or not os.path.isfile(path):
        return None
    try:
        with open(path, "r", encoding="utf-8") as f:
            raw = f.read().lstrip("\ufeff")
        s = raw.strip()
        if s.startswith("{"):
            data = json.loads(raw)
            return {str(k): str(v) for k, v in data.items()} or None
        m = {}
        for ln in raw.splitlines():
            line = ln.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                k, v = line.split("=", 1)
            elif "\t" in line:
                k, v = line.split("\t", 1)
            else:
                continue
            k, v = k.strip(), v.strip()
            if k and v:
                m[k] = v
        return m or None
    except Exception as e:
        print("  [警告] 无法解析自定义映射文件 %s：%s（已忽略）" % (path, e), file=sys.stderr)
        return None


def _make_skip_item(path, src_root, category, reason, action=None, warning=None):
    """构造结构化未处理项（skip manifest 条目）。"""
    act, warn = PREPROCESS_ACTIONS.get(category, (action, warning))
    if action is not None:
        act = action
    if warning is not None:
        warn = warning
    return {
        "path": path,
        "rel": os.path.relpath(path, src_root),
        "category": category,
        "reason": reason,
        "action": act,
        "warning": warn,
    }


def _skip(skipped, path, src_root, category, reason, action=None, warning=None, extra=None):
    """向 skipped 列表追加结构化未处理项（dict）。skipped 为 None 时跳过。
    分类与动作/警告来自 PREPROCESS_ACTIONS，允许调用处显式覆盖。
    extra 为可选 dict，合并进该项（如公开声明的 size/sha256）。"""
    if skipped is None:
        return
    item = _make_skip_item(path, src_root, category, reason, action, warning)
    if extra:
        item.update(extra)
    skipped.append(item)


def _file_meta(path):
    """计算文件 size/sha256（豁免留痕用；失败返回空 dict，绝不中断）。"""
    meta = {}
    try:
        meta["size"] = os.path.getsize(path)
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(1 << 16), b""):
                h.update(chunk)
        meta["sha256"] = h.hexdigest()
    except Exception:
        pass
    return meta


def _record_exempt(skipped, path, src_root, category, reason, warning):
    """记录豁免项（public_declared / local_only），附带 size/sha256（绝不静默）。"""
    _skip(skipped, path, src_root, category, reason,
          warning=warning, extra=_file_meta(path))


def _record_public_declared(skipped, path, src_root, reason):
    """记录「公开声明豁免」项（语义化入口）。"""
    _record_exempt(skipped, path, src_root, "public_declared", reason, PUBLIC_DECL_WARNING)


def _record_local_only(skipped, path, src_root, reason):
    """记录「本地处理·无需外发豁免」项（语义化入口）。"""
    _record_exempt(skipped, path, src_root, "local_only", reason, LOCAL_ONLY_WARNING)


def _copy_to_local_only(path, src_root, local_out_root):
    """工作区模式：把「仅本地处理」的未脱敏文件原样复制到隔离目录
    （07_本地处理不外发/），与可上云目录 03_脱敏副本/ 物理隔离（不得一起存放）。
    保持相对路径防同名覆盖；返回目标路径，失败返回 None。"""
    rel = os.path.relpath(path, src_root)
    dst = os.path.join(local_out_root, rel)
    os.makedirs(os.path.dirname(dst) or local_out_root, exist_ok=True)
    try:
        shutil.copy2(path, dst)
        return dst
    except Exception:
        return None


def _mk_skip(path, src_root, category, reason):
    """构造结构化未处理项（不追加；供 _classify_file 等返回）。"""
    return _make_skip_item(path, src_root, category, reason)


def _classify_file(path, src_root):
    """本地预处理分类：返回 dict（与 _skip 同构）或 None（treatable，无需预处理）。
    仅做轻量探测（图片不抽取；PDF 探测 /Encrypt + 试抽取；Office/db 试抽取），
    用于生成统一的“预处理清单”，不实际脱敏。"""
    ext = os.path.splitext(path)[1].lower()
    if ext in (TEXT_EXTS | CODE_EXTS):
        return None
    if ext in IMAGE_EXTS:
        return _mk_skip(path, src_root, "image",
                        "图片文件（无文本层，脚本不做 OCR）")
    if ext in EXTRACT_EXTS:
        if ext == ".pdf" and _pdf_has_encrypt(path):
            return _mk_skip(path, src_root, "encrypted",
                            "加密 PDF（含 /Encrypt 标记）")
        text, missing, reason = _extract_text(ext, path)
        if missing:
            return _mk_skip(path, src_root, "no_lib",
                            "缺少抽取库 %s" % missing)
        if text is None:
            if reason == "encrypted":
                return _mk_skip(path, src_root, "encrypted",
                                "加密文档（抽取异常/含保护）")
            if reason == "image_only":
                return _mk_skip(path, src_root, "image_pdf",
                                "纯图片型/扫描件 PDF（无文本层）")
            return _mk_skip(path, src_root, "error",
                            "抽取失败（%s）；可能已加密或文件损坏" % reason)
        if not text.strip():
            return _mk_skip(path, src_root, "empty",
                            "文档未提取到文本（可能为空或图片型）")
        return None  # 抽取成功，可直接脱敏
    return _mk_skip(path, src_root, "unsupported",
                    "未知/二进制格式（%s），若含敏感文本需先转文本" % ext)


# ---------------------------------------------------------------------------
# 本地预处理：自动解密（pikepdf / msoffcrypto-tool）+ 自动 OCR（rapidocr，纯本地离线）
# ---------------------------------------------------------------------------
def _decrypt_pdf(path, password, out_path):
    """用 pikepdf（封装 QPDF，C 实现、稳健）解密 PDF 并写出未脱敏的解密副本。
    返回 True=成功 / False=密码错误或解密失败 / None=缺少 pikepdf 库。
    pikepdf 保存时不设加密，故产出为未加密 PDF，可继续被分类/抽取/OCR 处理。"""
    try:
        import pikepdf
    except ImportError:
        return None
    try:
        with pikepdf.open(path, password=password) as src:
            src.save(out_path)
        return True
    except Exception:
        return False


def _decrypt_office(path, password, out_path):
    """用 msoffcrypto-tool 尝试解密 Office 文档并写出未脱敏副本。
    返回 True=成功 / False=失败 / None=缺少 msoffcrypto-tool 库。"""
    try:
        import msoffcrypto
    except ImportError:
        return None
    try:
        with open(path, "rb") as f:
            of = msoffcrypto.OfficeFile(f)
            of.load_key(password=password)
            with open(out_path, "wb") as o:
                of.decrypt(o)
        return True
    except Exception:
        return False


def _try_decrypt(path, passwords, out_path):
    """依次尝试候选密码解密；返回 (out_path, None) 或 (None, err)。"""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        dec, lib = _decrypt_pdf, "pikepdf"
    elif ext in (".docx", ".xlsx", ".pptx"):
        dec, lib = _decrypt_office, "msoffcrypto-tool"
    else:
        return None, "不支持自动解密的格式：%s（请手动解密后纳入）" % ext
    if not passwords:
        return None, "未提供候选密码（--passwords-file）"
    for pw in passwords:
        r = dec(path, pw, out_path)
        if r is True:
            return out_path, None
        if r is None:
            return None, ("缺少解密库 %s，无法自动解密：请在本环境安装后重跑"
                          "（uv add %s 或 pip install %s）" % (lib, lib, lib))
    return None, "所有候选密码均解密失败（密码错误或文件非标准加密）"


# ---------------------------------------------------------------------------
# 本地 OCR：rapidocr（纯 pip，模型随 wheel 捆绑，完全离线）
# ---------------------------------------------------------------------------
_OCR_ENGINES = {}
_OCR_ENGINE_ERR = None

# OCR 图像归一化与 det limit_side_len 自适应（精度优先，均经实测标定）：
# - 实测默认 limit=736 会把整页扫描件（~2100px）压缩到文字 ~15px → 识别碎片化；
#   而 ~2500px+ 区间 det 模型亦不稳定（实测 2527px 碎片化）；1000–2100px 为稳定区。
# - 故统一把图像归一化到最长边 ≤ _OCR_TARGET_PX（LANCZOS 降采样），
#   limit_side_len = clamp(最长边, 736, 2000)：不缩放为最佳，小图下限 736 略放大助识别。
_OCR_TARGET_PX = 2000
_OCR_LIMIT_MIN = 736
_OCR_LIMIT_MAX = 2000


def _ocr_limit_for(w, h):
    return min(max(max(w, h), _OCR_LIMIT_MIN), _OCR_LIMIT_MAX)


def _normalize_for_ocr(pil_img):
    """最长边超过目标像素则 LANCZOS 降采样到目标（避开 det 不稳定区间）。"""
    from PIL import Image
    w, h = pil_img.size
    m = max(w, h)
    if m > _OCR_TARGET_PX:
        r = _OCR_TARGET_PX / m
        pil_img = pil_img.resize((max(1, int(w * r)), max(1, int(h * r))),
                                 Image.LANCZOS)
    return pil_img


def _get_ocr_engine(limit_side_len):
    """按 limit_side_len 惰性初始化 rapidocr 引擎（进程内按配置缓存复用）。
    返回 (engine, None) 或 (None, err)。首次加载模型需数秒，之后极快。"""
    global _OCR_ENGINE_ERR
    if limit_side_len in _OCR_ENGINES:
        return _OCR_ENGINES[limit_side_len], None
    if _OCR_ENGINE_ERR is not None:
        return None, _OCR_ENGINE_ERR
    try:
        from rapidocr import RapidOCR
        engine = RapidOCR(params={"Det.limit_side_len": int(limit_side_len)})
        _OCR_ENGINES[limit_side_len] = engine
        return engine, None
    except ImportError:
        _OCR_ENGINE_ERR = "未安装 rapidocr/onnxruntime：" + OCR_NOTE
    except Exception as e:
        _OCR_ENGINE_ERR = "rapidocr 初始化失败：%s" % e
    return None, _OCR_ENGINE_ERR


def _ocr_result_texts(res):
    """从 rapidocr 返回结果中提取全部文本行（兼容 v3 结果对象与旧列表格式；
    未检测到文本时 txts 为 None，须按空处理）。"""
    if res is None:
        return []
    txts = getattr(res, "txts", None)
    if txts is not None:
        return [str(t) for t in txts if t and str(t).strip()]
    # 兼容 [[box, text, score], ...] 旧格式
    out = []
    try:
        for item in res:
            if isinstance(item, (list, tuple)) and len(item) >= 2 and item[1]:
                out.append(str(item[1]))
    except TypeError:
        pass
    return out


def _ocr_one_image(pil_img):
    """识别单张 PIL 图片（先归一化到稳定像素区）；返回 (text, None) 或 (None, err)。"""
    pil_img = _normalize_for_ocr(pil_img)
    engine, err = _get_ocr_engine(_ocr_limit_for(*pil_img.size))
    if engine is None:
        return None, err
    try:
        res = engine(pil_img)
    except Exception as e:
        return None, "rapidocr 识别失败：%s" % e
    return "\n".join(_ocr_result_texts(res)), None


def _ocr_image(path):
    """rapidocr 识别单张图片文件；返回 (text, None) 或 (None, err)。完全离线。"""
    try:
        from PIL import Image
    except ImportError:
        return None, "未安装 Pillow（rapidocr 依赖链应自带）"
    try:
        with Image.open(path) as im:
            return _ocr_one_image(im.convert("RGB"))
    except Exception as e:
        return None, "图片读取失败：%s" % e


def _ocr_pdf(path):
    """图片型 PDF：pypdfium2 逐页渲染为图像 → rapidocr 逐页识别。
    返回 (text, None) 或 (None, err)。完全离线。"""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return None, "未安装 pypdfium2（PDF 页面渲染必需）"
    texts = []
    try:
        doc = pdfium.PdfDocument(path)
        try:
            for i in range(len(doc)):
                page = doc[i]
                # 自适应渲染：让最长边≈目标像素（小页超采样提精度，大页控耗时），
                # 再由 _ocr_one_image 归一化兜底，避开 det 不稳定区间
                w_pt, h_pt = page.get_size()
                scale = _OCR_TARGET_PX / max(w_pt, h_pt) if max(w_pt, h_pt) else 2.0
                scale = min(max(scale, 1.0), 4.0)
                with page.render(scale=scale).to_pil() as pil_img:
                    text, err = _ocr_one_image(pil_img.convert("RGB"))
                if err is not None:
                    return None, err
                if text:
                    texts.append(text)
        finally:
            doc.close()
    except Exception as e:
        return None, "PDF 渲染/OCR 失败：%s" % e
    return "\n".join(texts), None


def _ocr_file(path):
    """按类型分流 OCR：PDF → 渲染后识别；图片 → 直接识别。"""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _ocr_pdf(path)
    return _ocr_image(path)


def _load_passwords(path):
    """读取候选密码文件：每行一个；或 JSON 列表 / {文件名: 密码} 字典。
    该密码文件本身含敏感信息，严禁外传。
    文件不存在/不可读时不崩溃：打印警告并按"无候选密码"处理。"""
    if not path:
        return []
    try:
        with open(path, "r", encoding="utf-8") as f:
            raw = f.read().strip()
    except OSError as e:
        print("  [警告] 密码文件读取失败（%s）：%s（按无候选密码处理）"
              % (path, e), file=sys.stderr)
        return []
    if not raw:
        return []
    if raw[0] in "[{":
        try:
            obj = json.loads(raw)
            if isinstance(obj, list):
                return [str(x) for x in obj]
            if isinstance(obj, dict):
                return [str(v) for v in obj.values()]
        except Exception:
            pass
    return [ln.strip() for ln in raw.splitlines() if ln.strip()]


def _write_preprocess_summary(path, items, exceptions, out_dir, manifest_path,
                              ready_list=None):
    """写出“本地预处理确认单”（满足：原文件/未脱敏副本保存外发清单、
    OCR 校对提醒、异常清单、确认后再 run 的闸门）。"""
    L = []
    L.append("# 本地预处理确认单（脱敏前必读 · 防割裂）\n")
    L.append("> **红线**：下方所有「原始文件」与「未脱敏副本（解密副本 / OCR 文本 / 直接可脱敏副本）」"
             "均 **本地留存、严禁外传**。\n"
             "> 仅当确认无误并完成脱敏后，脱敏副本方可上云。\n")
    # 一、保存 / 外发情况
    L.append("## 一、原始文件与未脱敏副本 · 保存 / 外发情况\n")
    L.append("| 原始文件 | 未脱敏副本（预处理产出） | 类型 | 保存位置 | 外发状态 |")
    L.append("| --- | --- | --- | --- | --- |")

    def _md(s):
        return str(s).replace("|", "\\|")

    for it in items:
        orig = it["original"]
        if it.get("preprocessed_copy"):
            copy = it["preprocessed_copy"]
        elif it.get("status") == "exception":
            copy = "（未产出，见异常清单）"
        else:
            copy = "（无需副本）"
        typ = it.get("preprocessed") or it["category"]
        L.append("| `%s` | `%s` | %s | 本地 | 🚫 禁止 |"
                 % (_md(orig), _md(copy), _md(typ)))
    for orig, copy in (ready_list or []):
        L.append("| `%s` | `%s` | ready（直接可脱敏·未脱敏副本） | 本地 | 🚫 禁止 |"
                 % (_md(orig), _md(copy)))
    L.append("")
    if any(it.get("preprocessed") == "decrypted+ocr" for it in items):
        L.append("> 说明：「decrypted+ocr」项的解密副本（ready/ 下的 PDF）无文本层，"
                 "其内容已由 ocr/ 下的 OCR 文本覆盖；后续 run 跳过该 PDF 属**正常现象**，"
                 "并非漏处理。\n")
    # 二、OCR 校对
    ocr_items = [it for it in items if it.get("needs_proofread")]
    L.append("## 二、OCR 结果 · 需用户逐项校对\n")
    if ocr_items:
        for it in ocr_items:
            L.append("- 校对文件 `%s`（对应原始 `%s`）：请逐字核对识别文本，"
                     "重点关注姓名 / 证件号 / 银行卡 / 金额等敏感字段；"
                     "确认无误后再脱敏。" % (it["preprocessed_copy"], it["original"]))
    else:
        L.append("- 本次无 OCR 产出，无需校对。")
    L.append("")
    # 三、异常清单
    L.append("## 三、预处理异常清单（须另行处理后发回 AI Agent）\n")
    L.append("| 文件 | 类别 | 异常原因 | 处理建议 |")
    L.append("| --- | --- | --- | --- |")
    if exceptions:
        for ex in exceptions:
            act = ex.get("action") or ("本地处理（解密 / 安装抽取库 / OCR / 转文本）后，"
                                       "将结果发回 AI Agent 重新纳入")
            L.append("| `%s` | %s | %s | %s |"
                     % (_md(ex["original"]), _md(ex["category"]),
                        _md(ex["exception"]), _md(act)))
    else:
        L.append("| —— | —— | 无异常 | —— |")
    L.append("")
    L.append("> 异常未处理前 **禁止** 执行 `run`；请用户本地处理后将结果交回 AI Agent。\n")
    # 四、确认与继续
    L.append("## 四、确认与继续\n")
    L.append("1. 确认：① 异常清单已全部处理；② OCR 文本已逐项校对无误。")
    L.append("2. 确认无误后执行（脚本会再次校验无异常才放行）：")
    L.append("   ```")
    L.append("   python desensitize.py run %s --preprocess-manifest %s --out <脱敏输出目录> --mode hybrid"
             % (out_dir, manifest_path))
    L.append("   ```")
    L.append("")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(L))


def cmd_preprocess(args):
    """本地预处理关卡（v2.4 起实际执行解密 / OCR——纯本地 rapidocr，
    模型内置、完全离线，无需任何外部服务）。

    流程：分类 → 自动解密（加密文档）/ 自动 OCR（图片、图片型 PDF）→ 产出
    “预处理确认单”（含原文件与未脱敏副本的保存/外发情况、OCR 校对提醒、异常清单），
    并在确认单中设置“经用户确认与校对无异常后再 run”的闸门。"""
    input_path = os.path.abspath(args.input)
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(input_path)
    ws = _ws_root(args)
    if ws:
        _ensure_workspace(ws)
        if args.out_dir == "./preprocessed":
            args.out_dir = os.path.join(ws, WS_UNDESEN)
        out_dir = os.path.abspath(args.out_dir)
        if not args.manifest:
            args.manifest = os.path.join(ws, "desensitize_preprocess.json")
        ready_dir = os.path.join(ws, WS_READY)
        ocr_dir = os.path.join(ws, WS_OCR)
        summary_path = os.path.join(ws, WS_PREPROC_SUMMARY)
    else:
        out_dir = os.path.abspath(args.out_dir)
        if not args.manifest:
            args.manifest = os.path.join(out_dir, "desensitize_preprocess.json")
        ready_dir = os.path.join(out_dir, "ready")
        ocr_dir = os.path.join(out_dir, "ocr")
        summary_path = os.path.join(out_dir, "desensitize_preprocess_summary.md")
    os.makedirs(ready_dir, exist_ok=True)
    os.makedirs(ocr_dir, exist_ok=True)

    passwords = _load_passwords(args.passwords_file)

    items, exceptions, treated_ready = [], [], []
    for p in iter_targets(input_path, args.recursive):
        c = _classify_file(p, src_root)
        if c is None:
            treated_ready.append(p)  # 直接可脱敏，拷入 ready（未脱敏副本）
            continue
        cat = c["category"]
        item = {
            "original": os.path.abspath(p),
            "rel": c["rel"],
            "category": cat,
            "external_send": "forbidden",
            "original_retained": True,
            "status": "done",
        }
        if cat == "encrypted":
            if args.no_auto:
                item["status"] = "needs_action"
                item["action"] = "本地解密（用户提供密码，仅本地操作）"
                item["warning"] = WARN_ENCRYPTED
                items.append(item)
                continue
            # 按相对路径保留目录结构，避免递归时不同子目录同名文件互相覆盖
            dec_path = os.path.join(ready_dir, c["rel"])
            os.makedirs(os.path.dirname(dec_path) or ready_dir, exist_ok=True)
            res, err = _try_decrypt(p, passwords, dec_path)
            if res is None:
                item["status"] = "exception"
                item["exception"] = err
                exceptions.append(item)
                items.append(item)
                continue
            item["preprocessed"] = "decrypted"
            item["preprocessed_copy"] = os.path.abspath(dec_path)
            item["warning"] = WARN_ENCRYPTED
            # 解密后可能仍是图片型（无文本层）→ 继续 OCR（本地 rapidocr，离线）
            c2 = _classify_file(dec_path, src_root)
            if c2 and c2["category"] in ("image_pdf", "image"):
                txt_path = os.path.join(ocr_dir, c["rel"] + ".ocr.txt")
                os.makedirs(os.path.dirname(txt_path) or ocr_dir, exist_ok=True)
                text, oerr = _ocr_file(dec_path)
                if text is not None:
                    _write_text_raw(txt_path, text)
                    item["preprocessed"] = "decrypted+ocr"
                    item["preprocessed_copy"] = os.path.abspath(txt_path)
                    item["needs_proofread"] = True
                else:
                    item["status"] = "exception"
                    item["exception"] = "解密成功但 OCR 失败：%s" % oerr
                    exceptions.append(item)
                    items.append(item)
                    continue
            items.append(item)
            continue
        if cat in ("image", "image_pdf"):
            if args.no_auto:
                item["status"] = "needs_action"
                item["action"] = "本地 rapidocr 识别为文本并落盘，提醒用户校对"
                item["warning"] = WARN_IMAGE
                items.append(item)
                continue
            txt_path = os.path.join(ocr_dir, c["rel"] + ".ocr.txt")
            os.makedirs(os.path.dirname(txt_path) or ocr_dir, exist_ok=True)
            text, oerr = _ocr_file(p)
            if text is not None:
                _write_text_raw(txt_path, text)
                item["preprocessed"] = "ocr"
                item["preprocessed_copy"] = os.path.abspath(txt_path)
                item["needs_proofread"] = True
                item["warning"] = WARN_IMAGE
                items.append(item)
                continue
            else:
                item["status"] = "exception"
                item["exception"] = oerr
                exceptions.append(item)
                items.append(item)
                continue
        # no_lib / error / empty / unsupported：列为异常（携带分类时的建议动作/警告）
        item["status"] = "exception"
        item["exception"] = c["reason"]
        item["action"] = c.get("action")
        item["warning"] = c.get("warning")
        exceptions.append(item)
        items.append(item)

    # 直接可脱敏的原始文件拷入 ready（作为未脱敏副本，禁止外发）；
    # 按相对路径保留目录结构，避免递归时同名文件互相覆盖
    for p in treated_ready:
        if os.path.basename(p) in META_SKIP:  # 不把流程自身元数据当业务文件拷贝
            continue
        dst = os.path.join(ready_dir, os.path.relpath(p, src_root))
        if os.path.abspath(p) != os.path.abspath(dst):
            os.makedirs(os.path.dirname(dst) or ready_dir, exist_ok=True)
            shutil.copy2(p, dst)

    # 写 manifest + 确认单
    manifest = {
        "time": datetime.now().isoformat(timespec="seconds"),
        "input": input_path,
        "out_dir": out_dir,
        "items": items,
        "exceptions": [{"original": e["original"], "category": e["category"],
                        "exception": e["exception"]} for e in exceptions],
        "has_exception": bool(exceptions),
        "note": "本确认单为脱敏前的本地预处理输出。encrypted/image/image_pdf 项已尝试自动"
                "解密/OCR；异常项须处理后再纳入。原文件与未脱敏副本严禁外传。",
    }
    with open(args.manifest, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    ready_list = [(os.path.abspath(p),
                   os.path.abspath(os.path.join(
                       ready_dir, os.path.relpath(p, src_root))))
                  for p in treated_ready
                  if os.path.basename(p) not in META_SKIP]
    _write_preprocess_summary(summary_path, items, exceptions, out_dir,
                              os.path.abspath(args.manifest), ready_list)

    # 控制台摘要
    print("本地预处理完成（自动解密 / rapidocr 本地 OCR，全程离线）%s：" %
          ("[递归]" if args.recursive else ""), file=sys.stderr)
    done = [it for it in items if it["status"] == "done"]
    need = [it for it in items if it["status"] == "needs_action"]
    print("  ✅ 已预处理（解密/OCR）：%d 项" % len(done), file=sys.stderr)
    print("  🚫 异常（须处理）：%d 项" % len(exceptions), file=sys.stderr)
    print("  📋 直接可脱敏（已拷入 ready/）：%d 项" % len(treated_ready), file=sys.stderr)
    if need:
        print("  ✎ 仅分类未处理（--no-auto）：%d 项" % len(need), file=sys.stderr)
    print("  预处理工作区：%s" % out_dir, file=sys.stderr)
    print("    - ready/ : 可直接脱敏文件（未脱敏副本·禁止外发）", file=sys.stderr)
    print("    - ocr/   : OCR 文本（未脱敏·需校对·禁止外发）", file=sys.stderr)
    print("  确认单（含保存/外发情况、OCR 校对、异常清单）：%s" % summary_path,
          file=sys.stderr)
    print("  清单 JSON：%s" % os.path.abspath(args.manifest), file=sys.stderr)
    if exceptions:
        print("  ⚠️ 存在 %d 项异常，请处理后再执行 run（run 也会拒绝异常未清）。"
              % len(exceptions), file=sys.stderr)
    else:
        print("  无异常：确认 OCR 校对无误后，执行 run（带 --preprocess-manifest 再次校验）。",
              file=sys.stderr)
    n_proof = sum(1 for it in items if it.get("needs_proofread"))
    if n_proof:
        print("  📝 有 %d 个 OCR 文件需逐项校对。" % n_proof, file=sys.stderr)
    if ws:
        _refresh_index(ws)
        print("  成果索引（一站式查阅全部产物）：%s" % os.path.join(ws, WS_INDEX_MD),
              file=sys.stderr)
        print("  下一步：校对 OCR 后执行 run（带 --workspace %s）" % ws, file=sys.stderr)


def cmd_scan(args):
    names = load_names(args.names)
    patterns = build_patterns(args.cn_enhance)
    input_path = os.path.abspath(args.input)
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(input_path)
    public_manifests = discover_public_manifests(input_path, getattr(args, "public_manifest", None))
    total = {}
    cleaning = {}
    skipped = []
    print("扫描命中（不生成文件）%s%s%s%s：" % (
        "[中文增强开启]" if args.cn_enhance else "",
        " [结构化列解析开启]" if args.tabular_names else "",
        " [公开声明豁免已启用]" if (args.assume_public or public_manifests) else "",
        " [本地处理豁免已启用]" if (args.local_only or args.local_paths) else ""))
    for p in iter_targets(input_path, args.recursive):
        process_file(p, "mask", "", src_root, {}, {}, names, total, do_write=False,
                     patterns=patterns, skipped=skipped,
                     public_manifests=public_manifests, assume_public=args.assume_public,
                     public_paths=args.public_paths, local_only=args.local_only,
                     local_paths=args.local_paths, cleaning=cleaning,
                     tabular_names=args.tabular_names)
    if total:
        print("\n汇总：", json.dumps(total, ensure_ascii=False))
    else:
        print("未发现已知敏感标识符。")
    if cleaning:
        print("\n⚠ 清洗建议（疑似未清洗数据形态，未自动脱敏）：%s" % json.dumps(cleaning, ensure_ascii=False))
        print(CLEANING_ADVICE_TEXT)
    _report_skipped(skipped, src_root)


def cmd_run(args):
    names = load_names(args.names)
    custom_map = load_custom_mapping(args.mapping)
    patterns = build_patterns(args.cn_enhance)
    input_path = os.path.abspath(args.input)
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(input_path)
    ws = _ws_root(args)
    if ws:
        _ensure_workspace(ws)
        if args.out == "./desensitized":
            args.out = os.path.join(ws, WS_DESEN)
        if args.keys == "./.desensitize_keys":
            args.keys = os.path.join(ws, WS_KEYS)
    out_root = args.out
    keys_dir = args.keys
    # 工作区模式：声明「仅本地处理」的文件原样复制到隔离目录（07_本地处理不外发/），
    # 与可上云目录 03_脱敏副本/ 物理隔离；非工作区模式不复制（仅留痕）。
    local_out_root = os.path.join(ws, WS_LOCAL_ONLY) if ws else None
    os.makedirs(out_root, exist_ok=True)

    # 预处理闸门：若提供了预处理清单且仍有异常，拒绝执行（满足“异常清完再继续”）。
    ppm_path = getattr(args, "preprocess_manifest", None)
    if ppm_path:
        # 预处理工作区（out_dir）内含 ready/ 与 ocr/ 子目录，强制递归扫描实际内容
        args.recursive = True
        try:
            with open(ppm_path, "r", encoding="utf-8") as f:
                ppm = json.load(f)
        except Exception as e:
            print("  [错误] 无法读取预处理清单 %s：%s" % (ppm_path, e), file=sys.stderr)
            sys.exit(2)
        if ppm.get("has_exception"):
            print("  [拒绝执行] 预处理清单仍存在未处理异常，禁止脱敏/上云。异常清单：",
                  file=sys.stderr)
            for ex in ppm.get("exceptions", []):
                print("    - %s  （%s）%s" % (ex.get("original"), ex.get("category"),
                      ex.get("exception")), file=sys.stderr)
            print("  请先本地处理上述异常（解密/安装库/OCR/转文本）后，将结果发回 AI Agent "
                  "重新预处理，再执行 run。", file=sys.stderr)
            sys.exit(1)
        print("  [闸门通过] 预处理清单无异常，继续执行脱敏。", file=sys.stderr)

    mapping = {}
    token_map = {}
    total = {}
    cleaning = {}
    skipped = []
    public_manifests = discover_public_manifests(input_path, getattr(args, "public_manifest", None))
    print("开始脱敏（模式=%s）%s%s%s%s..." % (
        args.mode,
        "[中文增强开启]" if args.cn_enhance else "",
        " [结构化列解析开启]" if args.tabular_names else "",
        " [公开声明豁免已启用]" if (args.assume_public or public_manifests) else "",
        " [本地处理豁免已启用]" if (args.local_only or args.local_paths) else ""))
    for p in iter_targets(input_path, args.recursive):
        process_file(p, args.mode, out_root, src_root, token_map, mapping, names, total,
                     do_write=True, patterns=patterns, skipped=skipped,
                     public_manifests=public_manifests, assume_public=args.assume_public,
                     public_paths=args.public_paths, local_only=args.local_only,
                     local_paths=args.local_paths, local_out_root=local_out_root,
                     cleaning=cleaning, custom_map=custom_map,
                     tabular_names=args.tabular_names)
    if cleaning:
        print("")
        print("  ⚠ 清洗建议（未清洗数据形态，未自动脱敏，请先清洗再脱敏）：%s" % cleaning)

    enc_info = encrypt_mapping(mapping, keys_dir, args.passphrase)
    # 唯一性检测：仅 mask 模式可能因“同型不同值”产生多对一歧义（影响恢复）。
    # token / hybrid 由唯一令牌保证 1:1，天然无歧义；redact 不可逆，恢复本就不适用——
    # 二者都不应进入碰撞告警（否则 redact 会因 [REDACTED_x] 全相同而误报告警）。
    if args.mode == "mask":
        collisions = find_collisions(mapping)
    else:
        collisions = {}
    if args.mode == "redact":
        safety = "irreversible"
    elif collisions:
        safety = "ambiguous"
    else:
        safety = "unique"
    report_collisions(collisions, args.mode)
    # 去掉映射表里的原始值明细，仅保留统计写入报告（原始值已在加密文件中）
    report = {
        "time": datetime.now().isoformat(timespec="seconds"),
        "input": os.path.abspath(args.input),
        "out": os.path.abspath(out_root),
        "mode": args.mode,
        "counts": total,
        "items_encrypted": sum(len(v) for v in mapping.values()),
        "restoration_safety": safety,
        "collision_risks": collisions,
        "mapping_encrypted": enc_info["mapping_file"],
        "key_file": enc_info.get("key_file"),
        "skipped_files": [it.get("rel") if isinstance(it, dict)
                          else os.path.relpath(it[0], src_root) for it in skipped],
        "skipped_detail": [
            {
                "file": it.get("rel") if isinstance(it, dict) else os.path.relpath(it[0], src_root),
                "category": it.get("category") if isinstance(it, dict) else "unsupported",
                "reason": it.get("reason") if isinstance(it, dict) else it[1],
                "action": it.get("action") if isinstance(it, dict) else None,
                "warning": it.get("warning") if isinstance(it, dict) else None,
                "size": it.get("size") if isinstance(it, dict) else None,
                "sha256": it.get("sha256") if isinstance(it, dict) else None,
            }
            for it in skipped
        ],
        "note": "原始敏感值仅存于加密映射表；脱敏副本位于 out 目录，可上云。"
                "restoration_safety=unique 表示可无歧义恢复（token / hybrid 模式恒为 unique）；"
                "=ambiguous 仅可能出现在 mask 模式（同型不同值碰撞），恢复可能混淆，"
                "建议改用 --mode hybrid（语义掩码+唯一令牌，兼顾字段语义与无歧义）或 --mode token；"
                "=irreversible 表示 redact 模式（不可逆，不可还原）。",
        # 清洗建议：疑似「未清洗/不规范」数据形态（带分隔符/短位手机、15 位旧身份证、
        # 订单号与银行卡区间重叠等）。这些形态未自动脱敏或可能误判，须先清洗再脱敏。
        "cleaning_advice": cleaning,
        "cleaning_advice_text": (CLEANING_ADVICE_TEXT if cleaning else ""),
    }
    # 报告位置：工作区模式下置于工作区根（位于脱敏副本目录之外，避免被一并上传）；
    # 非工作区模式保持原行为（取 out 的父目录）。
    if ws:
        report_path = os.path.join(ws, "desensitize_report.json")
    else:
        report_path = os.path.join(out_root, "..", "desensitize_report.json")
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    print("\n完成。")
    if ws:
        print("  脱敏副本目录 : %s" % os.path.abspath(out_root))
        print("  加密映射表   : %s" % WS_KEYS_MASK)
    else:
        # 非工作区模式：标注中文别名，便于理解与工作区目录的对应关系
        print("  脱敏副本目录 : %s  （= 工作区模式的 03_脱敏副本/）" % os.path.abspath(out_root))
        print("  加密映射表   : %s  （= 工作区模式的 04_映射表_保密/）" % WS_KEYS_MASK)
    if enc_info.get("key_file"):
        print("  密钥文件     : %s" % WS_KEYS_MASK)
    else:
        print("  密钥         : 由 --passphrase 派生（salt 已存于 keys 目录）")
    print("  命中统计     : %s" % json.dumps(total, ensure_ascii=False))
    print("  恢复安全性   : %s" % safety +
          ("（存在 %d 处多对一碰撞，恢复可能混淆）" % len(collisions)
           if safety == "ambiguous" else ""))
    print("  报告         : %s" % os.path.abspath(report_path))
    _report_skipped(skipped, src_root)

    # 工作区模式：生成上云前自检报告 + 刷新成果索引
    if ws:
        ocr_dir = os.path.join(ws, WS_OCR)
        sc = _write_selfcheck_report(ws, report, ocr_dir)
        _refresh_index(ws)
        print("")
        print("  📋 工作区成果索引 : %s" % os.path.join(ws, WS_INDEX_MD), file=sys.stderr)
        print("  🔒 上云前自检报告 : %s" % sc, file=sys.stderr)
        print("  → 上云前请逐项校对 OCR 副本 / 脱敏副本 / 本报告，确认后回复「确认上云」。",
              file=sys.stderr)


def _load_mapping(keys_dir, file=None, key=None, passphrase=None):
    """解密并返回映射表（data, enc_path）。供 decrypt / restore 复用。"""
    if file:
        enc_path = file
    else:
        encs = sorted(glob.glob(os.path.join(keys_dir, "mapping_*.json.enc")))
        if not encs:
            raise FileNotFoundError("未在 %s 找到 mapping_*.json.enc" % keys_dir)
        enc_path = encs[-1]
    if passphrase:
        salts = sorted(glob.glob(os.path.join(keys_dir, "salt_*.bin")))
        if not salts:
            raise ValueError("未找到 salt 文件，无法用口令解密")
        salt = open(salts[-1], "rb").read()
        k = derive_key_from_passphrase(passphrase, salt)
    else:
        keys = sorted(glob.glob(os.path.join(keys_dir, "desensitize_key_*.key")))
        if not keys:
            raise ValueError("未在 %s 找到 desensitize_key_*.key（也未提供 --passphrase）"
                             % keys_dir)
        k = open(key or keys[-1], "rb").read()
    fernet = make_fernet(k)
    data = json.loads(fernet.decrypt(open(enc_path, "rb").read()).decode("utf-8"))
    return data, enc_path


def cmd_decrypt(args):
    """解密映射表，供本地复核可逆性（无需上云）。"""
    ws = _ws_root(args)
    if ws and not args.out:
        args.out = os.path.join(ws, "映射表明细.json")
    try:
        data, enc_path = _load_mapping(args.keys, args.file, args.key, args.passphrase)
    except (FileNotFoundError, ValueError) as e:
        print(str(e), file=sys.stderr)
        return
    except Exception as e:
        print("解密失败：%s" % e, file=sys.stderr)
        return

    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print("已解密映射表写入：%s" % os.path.abspath(args.out))
        if ws:
            _refresh_index(ws)
        return

    total = sum(len(v) for v in data.values())
    # 过滤 redact 模式产生的 [REDACTED_*]：不可逆、本不适用恢复，不应作为歧义告警
    collisions = {r: v for r, v in find_collisions(data).items()
                  if not str(r).startswith("[REDACTED_")}
    print("映射表文件：%s" % enc_path)
    print("涉及文件 %d 个，命中条目 %d 条：" % (len(data), total))
    for f, items in data.items():
        print("  %-40s %d 条" % (f, len(items)))
    if collisions:
        print("\n[提醒] 该映射表存在 %d 处“同一脱敏值→多个原始值”的歧义，恢复时可能混淆："
              % len(collisions), file=sys.stderr)
        for r, origs in collisions.items():
            preview = " / ".join(origs[:8]) + (" …" if len(origs) > 8 else "")
            print("    %s  <-  %s" % (r, preview), file=sys.stderr)
        print("  → 如需无歧义恢复，请对该数据改用 --mode hybrid 或 --mode token 重新脱敏。",
              file=sys.stderr)
    else:
        print("\n[OK] 未检测到多对一歧义，原始值可唯一还原"
              "（hybrid / token 模式，或碰撞为空的 mask 模式）。")
    print("\n（加 --out <文件.json> 可导出完整 original↔replacement 明细，用于本地复核）")
    if ws:
        _refresh_index(ws)


def cmd_restore(args):
    """用加密映射表把脱敏副本回填为含原值的内部文档（映射表不离本地）。"""
    ws = _ws_root(args)
    if ws:
        _ensure_workspace(ws)
        if args.input == "./desensitized":
            args.input = os.path.join(ws, WS_DESEN)
        if args.out == "./restored":
            args.out = os.path.join(ws, WS_RESTORE)
    try:
        data, enc_path = _load_mapping(args.keys, args.file, args.key, args.passphrase)
    except (FileNotFoundError, ValueError) as e:
        print(str(e), file=sys.stderr)
        return
    except Exception as e:
        print("解密失败：%s" % e, file=sys.stderr)
        return

    input_root = args.input
    out_root = args.out
    os.makedirs(out_root, exist_ok=True)
    types = set(t.strip() for t in (args.types or "").split(",") if t.strip())

    total_restored = 0
    total_collision = 0
    total_skipped = 0
    for rel_path, items in data.items():
        # 定位脱敏副本：文本/代码按原扩展名；抽取型(.db/.pdf…)实际输出为 .txt
        cand = os.path.join(input_root, rel_path)
        suffix = ""
        if not os.path.isfile(cand):
            cand2 = cand + ".txt"
            if os.path.isfile(cand2):
                cand = cand2
                suffix = ".txt"
            else:
                print("  [跳过] 未找到脱敏副本：%s" % rel_path, file=sys.stderr)
                total_skipped += 1
                continue
        content = _read_text_plain(cand)

        # 构建 replacement -> original；检测多对一歧义（mask 同型/*** 密钥）
        repl_map = {}
        collisions = {}
        for it in items:
            if types and it.get("type") not in types:
                continue
            r = it["replacement"]
            o = it["original"]
            if r in repl_map and repl_map[r] != o:
                collisions.setdefault(r, set([repl_map[r]])).add(o)
            else:
                repl_map[r] = o
        for r in collisions:
            repl_map.pop(r, None)
        if collisions:
            total_collision += len(collisions)
            for r, origs in collisions.items():
                preview = " / ".join(sorted(origs))
                print("  [警告] 替换串 %r 对应多个原始值（%s），无法唯一还原，已跳过"
                      % (r, preview[:120]), file=sys.stderr)

        # 单次遍历替换（最长优先，避免子串互相干扰）；令牌模式请确保数据不含同名串
        if repl_map:
            escaped = sorted(((re.escape(r), r) for r in repl_map),
                             key=lambda x: -len(x[0]))
            pat = re.compile("|".join(e[0] for e in escaped))
            content = pat.sub(lambda m: repl_map[m.group(0)], content)
            total_restored += len(repl_map)

        out_path = os.path.join(out_root, rel_path) + suffix
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        _write_text_raw(out_path, content)
        print("  回填 -> %s （%d 处）" % (os.path.relpath(out_path, out_root), len(repl_map)))

    print("\n完成。回填目录: %s" % os.path.abspath(out_root))
    print("  回填条目: %d | 因碰撞跳过: %d | 未找到副本: %d" % (
        total_restored, total_collision, total_skipped))
    print("  （映射表 %s 已用于回填；回填产物恢复为含原值的本地内部文档，请按需谨慎保管，"
          "勿随脱敏副本一同外传）" % enc_path)
    if ws:
        _refresh_index(ws)


def cmd_audit(args):
    """基于 run 生成的 desensitize_report.json 自动产出审计文档（九节，
    对齐 SKILL.md 的 12 项上云前自查清单）。"""
    ws = _ws_root(args)
    if ws:
        _ensure_workspace(ws)
        if args.report == "./desensitize_report.json":
            args.report = os.path.join(ws, "desensitize_report.json")
        if args.out == "./desensitize_audit.md":
            args.out = os.path.join(ws, WS_AUDIT)
    report_path = args.report
    if not os.path.isfile(report_path):
        print("未找到报告文件：%s" % report_path, file=sys.stderr)
        return
    rep = json.load(open(report_path, encoding="utf-8"))
    time = rep.get("time", "")
    inp = rep.get("input", "")
    out = rep.get("out", "")
    mode = rep.get("mode", "")
    counts = rep.get("counts", {})
    safety = rep.get("restoration_safety", "")
    collisions = rep.get("collision_risks", {})
    mapping_file = rep.get("mapping_encrypted", "")
    key_file = rep.get("key_file") or "（由 --passphrase 派生，salt 存于 keys 目录）"
    skipped = rep.get("skipped_files", [])
    skipped_detail = rep.get("skipped_detail", [])
    items = rep.get("items_encrypted", 0)

    # 简化分级：直接标识符 / 密钥类 → 高；其余 → 中
    HIGH = {"id_card", "phone", "bank_card", "email", "ip", "plate",
            "passport", "name", "cn_name", "secret"}
    levels = {}
    for k, v in counts.items():
        lvl = "高" if k in HIGH else "中"
        levels.setdefault(lvl, {})[k] = v

    L = []
    L.append("# 脱敏审计记录（由 desensitize.py audit 自动生成）")
    L.append("")
    # ── 审计摘要（置顶，供非专业人员一屏内掌握结论；细节见后九节）──
    _risk_summary = "低" if (safety == "unique" and not skipped_detail) else "中"
    _upload_scope = "仅 `03_脱敏副本/`（已脱敏，可上云）；原始文件 / 未脱敏副本 / 映射表均留本地"
    _compliance = "已脱敏、可上云，原始与映射表留本地分离" if _risk_summary == "低" else "存在歧义或未处理文件，须人工复核后再上云"
    L.append("## 审计摘要")
    L.append("- 风险自评：**%s**" % _risk_summary)
    L.append("- 外发范围：%s" % _upload_scope)
    L.append("- 结论：%s" % _compliance)
    L.append("- 脱敏方式：%s（%s）" % (mode, safety))
    L.append("")
    L.append("> 完整审计见下方九节；上云前仍须逐项确认 12 项自查清单，禁止「一键脱敏即上云」。")
    L.append("")
    L.append("---")
    L.append("")
    L.append("- 生成时间：%s" % datetime.now().isoformat(timespec="seconds"))
    L.append("- 任务时间：%s" % time)
    L.append("- 任务描述：基于 `run` 报告自动生成；业务背景请人工补充")
    L.append("")
    L.append("## 一、数据盘点")
    L.append("- 原始数据：%s" % inp)
    L.append("- 脱敏条目（已加密入映射表）：%s 条" % items)
    L.append("- 脱敏副本目录：%s" % out)
    L.append("")
    L.append("## 二、敏感字段与级别")
    if levels:
        for lvl in ("高", "中", "低"):
            if lvl in levels:
                detail = "，".join("%s×%d" % (k, v)
                                  for k, v in sorted(levels[lvl].items()))
                L.append("- %s风险：%s" % (lvl, detail))
    else:
        L.append("- （报告未记录命中，或仅 redact 模式）")
    L.append("")
    L.append("## 三、脱敏方法")
    L.append("- 模式：%s" % mode)
    L.append("- 恢复安全性：%s%s" % (
        safety, "（存在碰撞，恢复可能混淆，建议改用 hybrid）" if safety == "ambiguous" else ""))
    if mode == "redact":
        L.append("- redact 不可逆，被抑制字段不可还原（精度：相关字段已删除）")
    else:
        L.append("- 数值精度：仅去标识、保留精确数值（未泛化/随机化），满足财务/审计红线")
    L.append("")
    L.append("### 脱敏前后对照示例（仅脱敏后形态，不含原值）")
    L.append("脱敏保留字段语义（可用于统计 / 分组），但无法还原个人身份：")
    L.append("")
    L.append("| 类型 | 脱敏后示例 | 说明 |")
    L.append("| --- | --- | --- |")
    L.append("| 手机号 | `138****8000` | 保留前 3 后 4，可用于区号 / 号段统计 |")
    L.append("| 身份证 | `4*****************` | 保留首位，其余掩码 |")
    L.append("| 姓名 | `张**` | 保留姓氏 / 首字 |")
    L.append("| 银行卡 | `6222****...` | 保留发卡行前缀（BIN） |")
    L.append("")
    L.append("> 精确的「原值 ↔ 令牌」对照仅存于加密映射表（重识别钥匙，留本地、不展示），不在此列出原值。")
    L.append("")
    L.append("## 四、密钥与映射表（重识别钥匙，须分离/加密/最小权限）")
    L.append("- 加密映射表：%s" % mapping_file)
    L.append("- 密钥文件：%s （权限 600）" % key_file)
    L.append("- 状态：原始文件与映射表留本地、未随副本上云（对应 SKILL.md 清单 1/4/5）")
    L.append("")
    L.append("## 五、上云内容")
    L.append("- 仅脱敏副本：%s" % out)
    L.append("")
    L.append("## 六、本地复核结论")
    if safety == "unique":
        L.append("- 令牌唯一、无歧义恢复：PASS（可用 `decrypt`/`restore` 还原且无混淆）")
    elif safety == "ambiguous":
        L.append("- 存在 %d 处多对一歧义，恢复可能混淆：" % len(collisions))
        for r, origs in collisions.items():
            L.append("    - %s <- %s" % (r, " / ".join(origs[:8])))
    else:
        L.append("- redact 不可逆，不适用恢复校验")
    L.append("")
    L.append("## 七、异常清单 / 未处理文件（须逐条人工处理，切勿让后续流程跳过）")
    if skipped_detail:
        for d in skipped_detail:
            L.append("- [%s] %s" % (d.get("category", "unsupported"), d.get("file", "")))
            L.append("    - 原因：%s" % d.get("reason", ""))
            if d.get("action"):
                L.append("    - 建议动作：%s" % d.get("action"))
            if d.get("warning"):
                L.append("    - %s" % d.get("warning"))
        L.append("> 上述文件未脱敏，上云前须按建议本地处理（解密 / OCR / 安装库 / 人工复核）"
                 "或确认不含敏感信息；原始文件与未脱敏文件均不得直接外传（SKILL.md 本地预处理关卡）。")
    elif skipped:
        for s in skipped:
            L.append("- %s" % s)
        L.append("> 上述文件未脱敏，上云前须转文本/OCR/解密或确认不含敏感信息（SKILL.md 红线）")
    else:
        L.append("- 无（全部目标文件已处理或被显式忽略）")
    L.append("")
    L.append("## 八、外泄风险自评")
    risk = "低" if (safety == "unique" and not skipped_detail) else "中"
    L.append("- 风险等级：%s" % risk)
    L.append("- 理由：脱敏副本可上云；原始/映射表留本地分离；%s"
             % ("存在歧义/未处理文件需关注" if risk == "中"
                else "无歧义恢复、无遗漏"))
    L.append("")
    L.append("## 九、操作人与待办")
    L.append("- 操作人：AI Agent（本地执行）")
    L.append("- 异常与待办：%s" % ("无" if not (collisions or skipped_detail)
                                    else "见第六/七节，需人工复核后上云"))
    L.append("")
    L.append("> 本审计由 `desensitize.py audit` 自动生成；SKILL.md 规定的 12 项上云前自查"
             "清单与人工复核仍须由操作人逐项确认，禁止“一键脱敏即上云”。")

    out_md = args.out
    os.makedirs(os.path.dirname(out_md) or ".", exist_ok=True)
    with open(out_md, "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")
    print("已生成审计文档：%s" % os.path.abspath(out_md))
    if ws:
        _refresh_index(ws)


def cmd_status(args):
    """回显工作区成果索引：一站式查看全部产物位置与“可上云/保密”状态。"""
    ws = _ws_root(args)
    if not ws or not os.path.isdir(ws):
        print("未找到有效工作区：%s" % ws, file=sys.stderr)
        print("用法：desensitize.py status --workspace <工作区目录>", file=sys.stderr)
        return
    md = os.path.join(ws, WS_INDEX_MD)
    if not os.path.isfile(md):
        _refresh_index(ws)
    print("=" * 60)
    print("脱敏工作区成果索引：%s" % os.path.abspath(ws))
    print("=" * 60)
    for e in json.load(open(os.path.join(ws, WS_INDEX_JSON), encoding="utf-8"))["entries"]:
        print("  [%-4s] %-22s %s" % (e["tag"], e["name"], e.get("display", e["path"])))
    print("-" * 60)
    print("仅 `03_脱敏副本/` 可上云；上云前请阅读 `05_上云前自检报告.md` 并完成四处校对。")
    print("详情见：%s" % os.path.abspath(md))


# ---------------------------------------------------------------------------
# guide：流程指引 + 闸门决策表（「一张图看懂流程」在代码里的单一来源）
# ---------------------------------------------------------------------------
GUIDE_TEXT = """\
信息脱敏上云 SOP · 快速指引（desensitize.py）

核心原则：原始文件留本地，只把脱敏副本上云；映射表（重识别钥匙）本地加密、绝不外发。

【闸门决策表】任务执行前先判定（任一项命中即按对应动作）
  ① 仅本地处理·无需外发 → 不脱敏（三条护栏 + 外发即失效 + 多次外发每次检测）
  ② 显式声明公开/已脱敏 → 不脱敏（可上云，留痕 + 警示，绝不静默）
  ③ 检出敏感           → run 脱敏（preprocess→run→自查→上云→审计）
  ④ 未检出             → 直接执行（零额外负担）

【命令链】推荐顺序
  preprocess   本地预处理：自动解密加密文档 + rapidocr 识别图片/图片型 PDF → 确认单
  run          脱敏：生成脱敏副本 + 加密映射表（--workspace 启用一站式中文工作区）
  status       回显工作区成果索引（可上云 / 保密 / 不外发标签）
  [上云]       仅上传 03_脱敏副本/（经 05_上云前自检报告 + 确认闸门）
  audit        基于报告自动生成审计文档（对齐 12 项自查清单）
  decrypt      本地解密映射表复核可逆性（无需上云）
  restore      用映射表把脱敏副本回填为含原值内部文档（生成工资单等实名交付物）

【7 个子命令】
  scan         仅扫描报告命中（不生成文件）
  preprocess   预处理关卡（解密 / OCR + 确认单）
  run          脱敏 + 加密映射表（核心，默认 hybrid 模式）
  status       回显工作区索引
  audit        自动审计文档
  decrypt      解密映射表（本地复核）
  restore      回填为实名内部文档

【快速开始】
  # 推荐：一站式工作区（中文目录 + 上云前自检报告）
  desensitize.py run ./data/ --workspace ./工作区 --mode hybrid
  desensitize.py status --workspace ./工作区

  # 极简：非工作区（默认 ./desensitized + ./.desensitize_keys）
  desensitize.py scan ./data/ --recursive
  desensitize.py run ./data/ --mode hybrid
"""


def cmd_guide(args):
    print(GUIDE_TEXT)


def build_parser():
    p = argparse.ArgumentParser(description="一键脱敏本地脚本（数据不出本机）")
    sub = p.add_subparsers(dest="cmd", required=True)

    sp = sub.add_parser("scan", help="扫描并报告敏感字段命中（不生成文件）")
    sp.add_argument("input", help="文件或目录")
    g = sp.add_argument_group("核心参数")
    g.add_argument("--recursive", action="store_true", help="递归处理目录")
    g = sp.add_argument_group("识别增强")
    g.add_argument("--names", help="已知姓名清单文件（每行一个）")
    g.add_argument("--cn-enhance", action="store_true",
                   help="中文识别增强：额外识别中文姓名/地址/机构名（本地正则，离线）")
    g.add_argument("--tabular-names", action="store_true",
                   help="结构化列解析：解析 CSV/TSV 首行列头，识别姓名列并脱敏英文姓名"
                        "（John Smith 等 Title Case 二词；本地正则、离线、零依赖）")
    g = sp.add_argument_group("豁免声明（用户显式声明才跳过脱敏，默认全脱敏）")
    g.add_argument("--assume-public", action="store_true",
                   help="整体声明公开/样例/已脱敏，全部跳过脱敏（仅留痕+警示，不脱敏不计数）")
    g.add_argument("--public-paths", nargs="+", default=None,
                   help="指定公开文件/文件夹路径（可多个）：文件夹=其下全部（递归）无需脱敏；"
                        "文件=该文件无需脱敏；仅本次任务生效（不落盘为永久设置）")
    g.add_argument("--public-manifest", default=None,
                   help="自定义公开声明伴随清单路径（.nodesens / desensitize_manifest.json）；"
                        "默认自动发现输入目录树内的清单（重复使用的高级选项）")
    g.add_argument("--local-only", action="store_true",
                   help="整体声明仅本地处理·无需外发，全部跳过脱敏（仅留痕+警示；绝不外发）")
    g.add_argument("--local-paths", nargs="+", default=None,
                   help="指定仅本地处理的文件/文件夹路径（可多个）：文件夹=其下全部（递归）；"
                        "内容保密但只本地处理、绝不外发；仅本次任务生效")
    sp.set_defaults(func=cmd_scan)

    rp = sub.add_parser("run", help="脱敏并生成加密映射表（核心命令，默认 hybrid 模式）")
    rp.add_argument("input", help="文件或目录")
    g = rp.add_argument_group("核心参数")
    g.add_argument("--out", default="./desensitized", help="脱敏副本输出目录（默认 ./desensitized，可上云）")
    g.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录（默认 ./.desensitize_keys，绝不外发）")
    g.add_argument("--mode", choices=["mask", "token", "redact", "hybrid"], default="hybrid",
                   help="脱敏模式（默认 hybrid）：mask=掩码(可逆) / token=令牌化(可逆,跨记录关联) / "
                        "hybrid=语义掩码+唯一令牌(可逆,无歧义,保留字段语义) / redact=抑制(不可逆)")
    g.add_argument("--recursive", action="store_true", help="递归处理目录")
    g.add_argument("--workspace", default=None,
                   help="统一成果中心（工作区）目录（推荐）：产物归拢一处、中文目录命名，"
                        "自动生成 成果索引 与 上云前自检报告（opt-in，不带则用非工作区默认路径）")
    g = rp.add_argument_group("识别增强")
    g.add_argument("--names", help="已知姓名清单文件（每行一个）")
    g.add_argument("--cn-enhance", action="store_true",
                   help="中文识别增强：额外识别中文姓名/地址/机构名（本地正则，离线）")
    g.add_argument("--tabular-names", action="store_true",
                   help="结构化列解析：解析 CSV/TSV 首行列头，识别姓名列并脱敏英文姓名"
                        "（John Smith 等 Title Case 二词；本地正则、离线、零依赖）")
    g = rp.add_argument_group("豁免声明（用户显式声明才跳过脱敏，默认全脱敏）")
    g.add_argument("--assume-public", action="store_true",
                   help="整体声明公开/样例/已脱敏，全部跳过脱敏（仅留痕+警示，不脱敏不计数）")
    g.add_argument("--public-paths", nargs="+", default=None,
                   help="指定公开文件/文件夹路径（可多个）：文件夹=其下全部（递归）无需脱敏；"
                        "文件=该文件无需脱敏；仅本次任务生效")
    g.add_argument("--public-manifest", default=None,
                   help="自定义公开声明伴随清单路径（.nodesens / desensitize_manifest.json）")
    g.add_argument("--local-only", action="store_true",
                   help="整体声明仅本地处理·无需外发，全部跳过脱敏（仅留痕+警示；绝不外发）")
    g.add_argument("--local-paths", nargs="+", default=None,
                   help="指定仅本地处理的文件/文件夹路径（可多个）：文件夹=其下全部（递归）；"
                        "内容保密但只本地处理、绝不外发；仅本次任务生效")
    g = rp.add_argument_group("高级")
    g.add_argument("--passphrase", default=None, help="用口令派生密钥（密钥不落盘）")
    g.add_argument("--preprocess-manifest", default=None,
                   help="预处理清单 JSON；若仍含异常则拒绝执行 run（异常清完再继续）")
    g.add_argument("--mapping", default=None,
                   help="用户自定义脱敏映射文件（原始值→替换值）：JSON 或文本每行 原始值=替换值；"
                        "命中优先用用户替换值；⚠ 该文件含原始值，严禁外传、仅本地使用，"
                        "命中结果随加密映射表本地保存、绝不外发")
    rp.set_defaults(func=cmd_run)

    dp = sub.add_parser("decrypt", help="解密映射表，供本地复核可逆性（不依赖上云）")
    dp.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录")
    dp.add_argument("--file", default=None, help="指定 mapping_*.json.enc（默认取目录中最新一个）")
    dp.add_argument("--key", default=None, help="指定密钥文件（默认取目录中最新 .key）")
    dp.add_argument("--passphrase", default=None, help="若用口令派生密钥，提供同一口令")
    dp.add_argument("--out", default=None, help="导出解密后的映射表 JSON 到该路径")
    dp.add_argument("--workspace", default=None,
                    help="统一成果中心（工作区）目录：未指定 --out 时，解密明细写入工作区并刷新成果索引")
    dp.set_defaults(func=cmd_decrypt)

    rp2 = sub.add_parser("restore", help="用映射表把脱敏副本回填为含原值的内部文档")
    rp2.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录")
    rp2.add_argument("--input", default="./desensitized",
                     help="脱敏副本目录（默认 ./desensitized；按映射表中的相对路径在该目录内定位副本）")
    rp2.add_argument("--out", default="./restored",
                     help="回填后内部文档输出目录（默认 ./restored）")
    rp2.add_argument("--file", default=None, help="指定 mapping_*.json.enc（默认取目录中最新一个）")
    rp2.add_argument("--key", default=None, help="指定密钥文件（默认取目录中最新 .key）")
    rp2.add_argument("--passphrase", default=None, help="若用口令派生密钥，提供同一口令")
    rp2.add_argument("--types", default=None,
                     help="仅回填指定类型，逗号分隔（如 name,id_card）；默认全部回填")
    rp2.add_argument("--workspace", default=None,
                    help="统一成果中心（工作区）目录：回填产物归入工作区并刷新成果索引")
    rp2.set_defaults(func=cmd_restore)

    ap = sub.add_parser("audit", help="基于 run 报告自动生成审计文档（九节，对齐 12 项自查清单）")
    ap.add_argument("--report", default="./desensitize_report.json",
                    help="run 生成的报告（默认 ./desensitize_report.json）")
    ap.add_argument("--out", default="./desensitize_audit.md",
                    help="审计文档输出路径（默认 ./desensitize_audit.md）")
    ap.add_argument("--workspace", default=None,
                    help="统一成果中心（工作区）目录：审计文档写入工作区并刷新成果索引")
    ap.set_defaults(func=cmd_audit)

    pp = sub.add_parser("preprocess",
                         help="本地预处理关卡（自动解密 / rapidocr 本地 OCR）：生成预处理确认单与异常清单")
    pp.add_argument("input", help="文件或目录")
    pp.add_argument("--recursive", action="store_true", help="递归处理目录")
    pp.add_argument("--out-dir", default="./preprocessed",
                    help="预处理工作区（默认 ./preprocessed，含 ready/ 与 ocr/ 子目录）")
    pp.add_argument("--manifest", default=None,
                    help="预处理清单 JSON（默认 <out-dir>/desensitize_preprocess.json）")
    pp.add_argument("--passwords-file", default=None,
                    help="解密候选密码文件：每行一个；或 JSON 列表 / {文件名:密码}。"
                         "该文件含敏感信息，严禁外传")
    pp.add_argument("--no-auto", action="store_true",
                    help="仅分类与提醒（旧行为），不实际解密/OCR")
    pp.add_argument("--workspace", default=None,
                    help="统一成果中心（工作区）目录：启用后未脱敏副本/OCR文本/确认单归入一处"
                         "并生成成果索引（opt-in，不带则保持原默认行为）")
    pp.set_defaults(func=cmd_preprocess)

    st = sub.add_parser("status", help="回显工作区成果索引（一站式查看全部产物位置与可上云状态）")
    st.add_argument("--workspace", required=True, help="工作区目录（由 preprocess/run 的 --workspace 指定）")
    st.set_defaults(func=cmd_status)

    gd = sub.add_parser("guide", help="打印流程指引与闸门决策表（一张图看懂流程）")
    gd.set_defaults(func=cmd_guide)
    return p


def main():
    args = build_parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
