#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
一键脱敏本地脚本（desensitize.py）

功能：在本地对文本 / 数据文件中的敏感标识符进行脱敏，生成脱敏副本，
并将「可逆映射表」加密保存到独立目录，确保原始文件不离开本机。

设计原则（对齐《敏感信息脱敏操作 SOP》）：
- 本地优先：脚本只在本地运行，原始文件与映射表永不离开本机。
- 最小必要：只处理识别到的敏感字段。
- 映射表分离 + AES 加密：脱敏副本可上云；映射表单独存放并加密（密钥不入库）。
- 中文优先：内置中文身份证 / 手机号 / 银行卡 / 邮箱 / 车牌 / 护照等正则。
- 中文增强（--cn-enhance）：本地正则识别中文姓名 / 地址 / 机构名，
  纯离线、无需安装额外模型。

依赖：
- 必需：cryptography（AES / Fernet）
- 中文增强：纯内置正则，无额外依赖（不依赖 Presidio 等需联网下载模型的方案）

用法：
  # 扫描（仅报告命中，不生成文件）
  uv run --project <tool_dir> python desensitize.py scan ./data/

  # 脱敏（默认 hybrid 模式），输出到 ./desensitized，映射表到 ./.desensitize_keys
  uv run --project <tool_dir> python desensitize.py run ./data/ \
      --out ./desensitized --keys ./.desensitize_keys

  # 令牌化（可逆，便于跨记录关联计数 / 分组）
  uv run --project <tool_dir> python desensitize.py run ./data/ --mode token

  # 语义掩码 + 唯一令牌（hybrid）：保留字段语义（张*/138****8000…）且可无歧义恢复
  uv run --project <tool_dir> python desensitize.py run ./data/ --mode hybrid --cn-enhance

  # 中文识别增强：额外识别中文姓名 / 地址 / 机构名
  uv run --project <tool_dir> python desensitize.py run ./data/ --cn-enhance
  uv run --project <tool_dir> python desensitize.py scan ./data/ --cn-enhance --recursive

  # 用口令派生密钥（密钥不落盘，靠口令恢复映射表）
  uv run --project <tool_dir> python desensitize.py run ./data/ --passphrase "***"

  # 本地复核：解密映射表，验证可逆性（原始值不出本机、无需上云）
  uv run --project <tool_dir> python desensitize.py decrypt --keys ./.desensitize_keys
  uv run --project <tool_dir> python desensitize.py decrypt --keys ./.desensitize_keys --out mapping_review.json

  # 回填：用映射表把脱敏副本还原为含原值的本地内部文档（映射表不离本地）
  uv run --project <tool_dir> python desensitize.py restore --keys ./.desensitize_keys \
      --input ./desensitized --out ./restored
  # 仅回填姓名与身份证（如工资条只需回写姓名）
  uv run --project <tool_dir> python desensitize.py restore --keys ./.desensitize_keys \
      --types name,id_card --out ./restored_names

  # 审计：基于 run 报告自动生成 11 项审计文档
  uv run --project <tool_dir> python desensitize.py audit \
      --report ./desensitize_report.json --out ./desensitize_audit.md
"""

import argparse
import glob
import json
import os
import re
import secrets
import sys
from datetime import datetime

from cryptography.fernet import Fernet
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC

# ---------------------------------------------------------------------------
# 1. 敏感字段正则（中文优先）
# ---------------------------------------------------------------------------
# 说明：身份证 / 手机号先处理，其数字会被替换为 '*'，后续数字类正则不会重复命中。
#
# 边界设计：原先用 \b 词边界，但 Python 正则里 CJK 也算"词字符"，
# 导致"数字紧挨中文"（如 手机138...、身份证110...）时 机1/证1 之间无词边界、
# 标识符漏检。现改用自定义边界：
#   _BOUND_L = 左侧不得是 ASCII 字母/数字（CJK、标点、空白、串首 均视为边界）
#   _BOUND_R = 右侧不得是 ASCII 字母/数字（CJK、标点、空白、串尾 均视为边界）
# 效果：① 数字紧邻中文可被正确识别（召回提升）；
#       ② 仍阻止"在更长字母/数字串内部"误匹配（如 a138... / 138…9 不误命中）。
_BOUND_L = r"(?<![0-9A-Za-z])"
_BOUND_R = r"(?![0-9A-Za-z])"

PATTERNS = {
    "id_card": re.compile(_BOUND_L + r"[1-9]\d{5}(?:18|19|20)\d{2}(?:0[1-9]|1[0-2])(?:0[1-9]|[12]\d|3[01])\d{3}[\dXx]" + _BOUND_R),
    "phone": re.compile(_BOUND_L + r"1[3-9]\d{9}" + _BOUND_R),
    "bank_card": re.compile(_BOUND_L + r"\d{16,19}" + _BOUND_R),
    "ip": re.compile(_BOUND_L + r"(?:(?:25[0-5]|2[0-4]\d|1?\d?\d)\.){3}(?:25[0-5]|2[0-4]\d|1?\d?\d)" + _BOUND_R),
    "email": re.compile(_BOUND_L + r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}" + _BOUND_R),
    "plate": re.compile(_BOUND_L + r"[京津沪渝冀豫云辽黑湘皖鲁新苏浙赣鄂桂甘晋蒙陕吉闽贵粤青藏川宁琼使领][A-Z][A-Z0-9]{5,6}" + _BOUND_R),
    "passport": re.compile(_BOUND_L + r"[EGDSH]\d{8}" + _BOUND_R + r"|" + _BOUND_L + r"[a-zA-Z]\d{9}" + _BOUND_R),
}

# ---------------------------------------------------------------------------
# 1.1 中文识别增强正则（--cn-enhance 时启用；纯本地，无需联网/模型）
# ---------------------------------------------------------------------------
# 常见复姓（须置于单姓之前优先匹配）
_CN_COMPOUND = ["欧阳", "司马", "上官", "诸葛", "东方", "皇甫", "令狐", "夏侯",
                "宇文", "慕容", "司徒", "端木", "拓跋", "轩辕", "公孙", "长孙",
                "鲜于", "耶律", "完颜", "万俟", "司空", "子车", "颛孙"]
# 常见单姓（百家姓 + 常见姓）
_CN_SINGLE = ("赵钱孙李周吴郑王冯陈褚卫蒋沈韩杨朱秦尤许何吕施张孔曹严华金魏陶姜"
              "戚谢邹喻柏水窦章云苏潘葛奚范彭郎鲁韦昌马苗凤花方俞任袁柳酆鲍史唐"
              "费廉岑薛雷贺倪汤滕殷罗毕郝邬安常乐于时傅皮卞齐康伍余元卜顾孟平黄"
              "和穆萧尹姚邵湛汪祁毛禹狄米贝明臧计伏成戴谈宋茅庞熊纪舒屈项祝董梁"
              "杜阮蓝闵席季麻强贾路娄危江童颜郭梅盛林刁钟徐邱骆高夏蔡田樊胡凌霍"
              "万柯卢莫房裘缪干解应宗丁宣贲邓郁单杭洪包诸左石崔吉钮龚程嵇邢滑裴"
              "陆荣翁荀羊於惠甄曲家封芮羿储靳汲邴糜松井段富巫乌焦巴弓牧隗山谷车"
              "侯宓蓬全郗班仰秋仲伊宫宁仇栾暴甘钭厉戎祖武符刘景詹束龙叶幸司韶"
              "郜黎蓟薄印宿白怀蒲台从鄂索咸籍赖卓蔺屠蒙池乔阴郁胥能苍双闻莘党"
              "翟谭贡劳逄姬申扶堵冉宰郦雍却璩桑桂濮牛寿通边扈燕冀郏浦尚农温别"
              "庄晏柴瞿阎充慕连茹习宦艾鱼容向古易慎戈廖庚终暨居衡步都耿满弘匡"
              "国文寇广禄阙东殴殳沃利蔚越夔隆师巩厍聂晁勾敖融冷訾辛阚那简饶空"
              "曾毋沙乜养鞠须丰巢关蒯相查后荆红游竺权逯盖益桓公")

# 姓名识别：姓氏 + 1~2 个汉字，依赖语境以抑制纯中文文本极高误报率。
# 分两级语境：
#  - 强语境（标签词 姓名/联系人/客户/借款人/担保人… 或 冒号 ：）：前缀已提供强上下文，
#    放宽后缀边界（仅按 姓氏+1~2 字 定长），提升召回（如"客户张三于..."可命中"张三"）。
#  - 弱语境（逗号/顿号/空格/括号/句号等标点前缀）：保留后缀边界（称谓或标点/结尾），
#    防止自由文本误报。前后边界同时覆盖 中英文标点（全角/半角），避免混排漏检。
# 中英文标点边界集（全角/半角一并包含）：空白、逗号、句号、冒号、分号、叹号、问号、顿号、
# 括号（全/半角）。供弱语境的前顾与后顾复用。
_CN_BOUND = r"\s,，.:：;；!！?？、()（）"
_CN_NAME_STRONG = re.compile(
    r"(?:(?<=姓名)|(?<=联系人)|(?<=客户)|(?<=经办人)|(?<=申请人)|"
    r"(?<=被谈话人)|(?<=被询问人)|(?<=借款人)|(?<=担保人)|"
    r"(?<=[:：]))"
    r"(?P<name>(?:%s|[%s])[一-龥]{1,2})"
    % ("|".join(_CN_COMPOUND), _CN_SINGLE)
)
_CN_NAME_WEAK = re.compile(
    r"(?<=[%s])"
    r"(?P<name>(?:%s|[%s])[一-龥]{1,2})"
    r"(?=(?:先生|女士|同学|老师|经理|总监|主任|教授|工|君|[%s]|$))"
    % (_CN_BOUND, "|".join(_CN_COMPOUND), _CN_SINGLE, _CN_BOUND)
)
CN_NAME_PATTERNS = [_CN_NAME_STRONG, _CN_NAME_WEAK]

# 地址识别：含 省/市/区/县 等行政区划 或 路/街/号/栋 等楼栋要素（二者至少其一）。
# 不强制"行政区划必须前置"，以覆盖"中关村大街1号"（无市/区）与"北京市海淀区"（无楼栋）等常见写法；
# 中英文标点在字符类之外，天然作为截断边界，混排不影响。
_CN_ADDRESS = re.compile(
    r"[一-龥0-9]{1,12}"
    r"(?:省|市|自治区|特别行政区|区|县|旗|盟|"
    r"路|街道|街|道|大道|巷|弄|号|栋|幢|单元|室|小区|花园|广场|大厦|公寓|中心|村|镇|乡)"
    r"[一-龥0-9A-Za-z\-]{0,16}"
)

# 机构识别：含 公司/银行/医院/大学 等组织后缀
_CN_ORG = re.compile(
    r"[一-龥]{2,}(?:公司|集团|银行|医院|大学|学院|学校|研究所|研究院|局|部|"
    r"委员会|协会|基金会|事务所|工厂|超市|店|分行|支行|证券|保险|基金|信托|"
    r"律所|合作社|中心|平台)"
)

CN_ENHANCE_PATTERNS = {
    "cn_name": CN_NAME_PATTERNS,
    "cn_address": _CN_ADDRESS,
    "cn_org": _CN_ORG,
}


def build_patterns(cn_enhance: bool):
    """返回生效的正则集合（cn_enhance 时并入中文增强）。"""
    pats = dict(PATTERNS)
    if cn_enhance:
        pats.update(CN_ENHANCE_PATTERNS)
    return pats


# 代码场景：密钥 / Token 类（仅对代码与配置文件生效）
# 值支持三种形态：单引号 '...'、双引号 "..."、或无引号（遇空白/引号截止）。
# 注意：引号型必须同时匹配「开头 + 结尾」引号，否则替换时结尾引号会残留（如 password='***'）。
SECRET_PATTERN = re.compile(
    r"""(?i)(['"]?)(?:api[_-]?key|secret|token|password|passwd|access[_-]?key|private[_-]?key|auth)['"]?\s*[:=]\s*(?P<val>'[^']{4,}'|"[^"]{4,}"|[^\s'"]{4,})"""
)

CODE_EXTS = {".py", ".js", ".ts", ".go", ".java", ".c", ".cpp", ".rb", ".php",
              ".yml", ".yaml", ".toml", ".env", ".ini", ".sh", ".sql"}
TEXT_EXTS = {".txt", ".md", ".csv", ".tsv", ".json", ".log", ".xml", ".html",
             ".htm", ".rst", ".eml"}
# 二进制/库文件：抽取文本后再脱敏（抽取库缺失时跳过并告警，不崩溃）
EXTRACT_EXTS = {".docx", ".xlsx", ".pptx", ".pdf", ".db"}
# 影像/图片：脚本不做 OCR，必须提醒用户本地 OCR 转文本后再纳入
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff",
              ".webp", ".heic", ".heif"}
# skip manifest 的明确可执行提醒（让用户知道“为什么没处理、下一步做什么”）
REMIND = {
    "encrypted": "加密文档：脚本无法处理加密文件，请先在本地解密后再纳入脱敏（解密后重跑本命令）",
    "image_only": "纯图片型/扫描件 PDF：pdfminer 仅能提取文本层，本文件未提取到任何文本，"
                  "请先在本地用 OCR 将图片转为文本（.txt/.md）后再纳入脱敏",
    "empty": "文档未提取到文本（可能为图片型扫描件或空文档）：若含敏感信息，"
             "请先 OCR 转文本后再纳入脱敏",
    "error": "文档读取/抽取失败（可能已加密或文件损坏）：请先解密或确认文件完整性后再纳入脱敏",
    "no_lib": "缺少抽取库，无法处理该二进制文档：请安装对应库（docx→python-docx，"
              "xlsx→openpyxl，pptx→python-pptx，pdf→pdfminer.six）后再处理",
}

# ---------------------------------------------------------------------------
# 2. 脱敏替换策略
# ---------------------------------------------------------------------------
def mask_value(kind: str, value: str) -> str:
    """按类型做掩码（可逆映射由调用方记录）。"""
    if kind == "phone" and len(value) == 11:
        return value[:3] + "****" + value[-4:]
    if kind == "id_card" and len(value) == 18:
        return value[:6] + "*" * 8 + value[-4:]
    if kind in ("bank_card",) and len(value) >= 12:
        return "*" * (len(value) - 4) + value[-4:]
    if kind == "email":
        local, _, domain = value.partition("@")
        keep = local[0] if local else "x"
        return keep + "***@" + domain
    if kind == "ip":
        parts = value.split(".")
        if len(parts) == 4:
            return parts[0] + "." + parts[1] + ".*.*"
        return "*.**.**.*"
    if kind == "plate" and len(value) >= 7:
        return value[:2] + "*" * (len(value) - 4) + value[-2:]
    if kind == "passport" and len(value) >= 8:
        return value[0] + "*" * (len(value) - 3) + value[-2:]
    if kind in ("cn_name", "name"):
        # 保留姓氏/首字，其余脱敏（原始值仍记录于加密映射表，可逆）
        return value[0] + "*" * (len(value) - 1)
    if kind == "cn_address":
        return "[地址]"
    if kind == "cn_org":
        return "[机构]"
    # 兜底：首尾各留 1，中间掩码
    if len(value) <= 2:
        return "*" * len(value)
    return value[0] + "*" * (len(value) - 2) + value[-1]


def redact_value(kind: str, value: str) -> str:
    return "[REDACTED_%s]" % kind.upper()


def token_value(token_map: dict, value: str) -> str:
    """为同一原始值生成稳定令牌（跨记录可关联）。"""
    if value in token_map:
        return token_map[value]
    tok = "T%04d" % (len(token_map) + 1)
    token_map[value] = tok
    return tok


# hybrid 模式分隔符：语义掩码 + 唯一令牌，形如  张*⟦T0001⟧
# 选用数学空心方括号（U+27E6 / U+27E7），在正常数据中极罕见，避免与正文混淆。
HYBRID_SEP_L, HYBRID_SEP_R = "⟦", "⟧"


# ---------------------------------------------------------------------------
# 3. 识别与脱敏核心
# ---------------------------------------------------------------------------
def desensitize_text(text: str, mode: str, token_map: dict, counts: dict,
                     names: set = None, patterns=None):
    """对单段文本做脱敏，返回 (脱敏后文本, 命中列表)。

    注意：必须一次性构建新串（re.sub + 回调），不可在 finditer 迭代中
    原地切片替换——否则变长替换会导致后续匹配偏移错位。
    替换串长度完全可变（可长可短，甚至远长于原文，如 hybrid 的“掩码⟦Txxxx⟧”）：
    re.sub 会按需动态拼接新串，无需、也无法在扫描阶段预分配固定长度，
    因此不存在“替换串长度不足”的失败模式——令牌超过 9999 个时自动扩展为
    T10000… 仍唯一且可被分隔符 ⟦ ⟧ 正确解析。
    """
    if patterns is None:
        patterns = PATTERNS
    hits = []
    result = text

    # 3.1 自定义姓名清单（若有）
    if names:
        name_cb = _repl_closure("name", mode, token_map, counts, hits)
        for nm in sorted(names, key=len, reverse=True):
            if not nm:
                continue
            result = re.sub(re.escape(nm), name_cb, result)

    # 3.2 顺序应用正则（先处理 ID/手机，避免银行卡重复命中）
    for kind in ["id_card", "phone", "bank_card", "ip", "email", "plate", "passport"]:
        pat = patterns[kind]
        cb = _repl_closure(kind, mode, token_map, counts, hits)
        result = pat.sub(lambda m, _cb=cb: _cb(m), result)

    # 3.3 中文增强（cn_name 为双模式列表，其余为单正则）
    if "cn_name" in patterns:
        for npat in patterns["cn_name"]:
            cb = _repl_closure("cn_name", mode, token_map, counts, hits)
            result = npat.sub(lambda m, _cb=cb: _cb(m), result)
    if "cn_address" in patterns:
        cb = _repl_closure("cn_address", mode, token_map, counts, hits)
        result = patterns["cn_address"].sub(lambda m, _cb=cb: _cb(m), result)
    if "cn_org" in patterns:
        cb = _repl_closure("cn_org", mode, token_map, counts, hits)
        result = patterns["cn_org"].sub(lambda m, _cb=cb: _cb(m), result)

    return result, hits


def _repl_closure(kind: str, mode: str, token_map: dict, counts: dict, hits: list):
    """生成脱敏回调：记录命中并返回替换串。"""
    def cb(m):
        seg = m.group(0)
        # 已被 ID/手机占位覆盖的纯数字段不应再当银行卡
        if kind == "bank_card" and set(seg) <= set("*"):
            return seg
        if mode == "redact":
            repl = redact_value(kind, seg)
        elif mode == "token":
            repl = token_value(token_map, seg)
        elif mode == "hybrid":
            # 语义掩码 + 唯一令牌：保留字段类型，同时靠令牌实现无歧义恢复
            repl = "%s%s%s%s" % (mask_value(kind, seg), HYBRID_SEP_L,
                                 token_value(token_map, seg), HYBRID_SEP_R)
        else:
            repl = mask_value(kind, seg)
        counts[kind] = counts.get(kind, 0) + 1
        hits.append({"type": kind, "original": seg, "replacement": repl})
        return repl
    return cb


def scan_text(text: str, names: set = None, patterns=None):
    """仅扫描，返回命中计数（复用脱敏的顺序逻辑，保证与 run 一致）。"""
    token_map = {}
    counts = {}
    desensitize_text(text, "mask", token_map, counts, names, patterns)
    return counts


def process_file(path: str, mode: str, out_root: str, src_root: str, token_map: dict,
                 mapping: dict, names: set, counts_total: dict, do_write: bool,
                 patterns=None, skipped=None):
    """处理单个文件。skipped 为可选 list，用于收集“未处理文件”（skip manifest）。"""
    try:
        ext = os.path.splitext(path)[1].lower()
        # 二进制/库文件：抽取文本后再脱敏（库缺失则跳过并告警）
        if ext in EXTRACT_EXTS:
            return _process_extracted(path, mode, out_root, src_root, token_map,
                                      mapping, names, counts_total, do_write,
                                      patterns, skipped)
        # 既非文本/代码也非可抽取类型：登记为未处理（不静默忽略）
        if ext not in (TEXT_EXTS | CODE_EXTS):
            if skipped is not None:
                if ext in IMAGE_EXTS:
                    skipped.append((path,
                        "影像/图片文件（%s）：本脚本不做 OCR，请先在本地用 OCR 将图片转为文本"
                        "（.txt/.md）后再纳入脱敏；原始图片不会被处理" % ext))
                else:
                    skipped.append((path,
                        "未支持的信息源类型（%s）：二进制/未知格式，若含敏感文本请先转为"
                        "文本文件（或 OCR/解密）后再纳入脱敏" % ext))
            return "skipped"
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
    except Exception as e:
        print("  [跳过] 无法读取 %s: %s" % (path, e), file=sys.stderr)
        if skipped is not None:
            skipped.append((path, "读取失败: %s" % e))
        return "error"

    ext = os.path.splitext(path)[1].lower()
    is_code = ext in CODE_EXTS
    rel_path = os.path.relpath(path, src_root)

    if not do_write:
        c = scan_text(text, names, patterns)
        if is_code:
            sec = len(SECRET_PATTERN.findall(text))
            if sec:
                c["secret"] = c.get("secret", 0) + sec
        if c:
            print("  %-40s %s" % (rel_path, c))
            for k, v in c.items():
                counts_total[k] = counts_total.get(k, 0) + v
        return "processed"

    # 代码文件额外处理密钥
    if is_code:
        def _sec_repl(m):
            val = m.group("val")
            # 剥离包裹引号，记录纯净原始值
            if len(val) >= 2 and val[0] in "'\"" and val[-1] == val[0]:
                inner = val[1:-1]
            else:
                inner = val
            counts_total["secret"] = counts_total.get("secret", 0) + 1
            if mode == "hybrid":
                repl = "***%s%s%s" % (HYBRID_SEP_L, token_value(token_map, inner),
                                      HYBRID_SEP_R)
            else:
                repl = "***"
            mapping.setdefault(rel_path, []).append(
                {"type": "secret", "original": inner, "replacement": repl})
            return m.group(0).replace(val, repl)
        text = SECRET_PATTERN.sub(_sec_repl, text)

    new_text, hits = desensitize_text(text, mode, token_map, counts_total, names, patterns)
    for h in hits:
        h["file"] = rel_path
        mapping.setdefault(rel_path, []).append(h)

    out_path = os.path.join(out_root, rel_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(new_text)
    return "processed"


# ---------------------------------------------------------------------------
# 3.1 二进制/库文件文本抽取（抽取后再脱敏；库缺失则跳过并告警）
# ---------------------------------------------------------------------------
def _lib_for(ext):
    return {".docx": "python-docx", ".xlsx": "openpyxl",
            ".pptx": "python-pptx", ".pdf": "pdfminer.six"}.get(ext, "未知库")


def _extract_sqlite(path):
    """用内置 sqlite3 把库表内容导出为文本（零依赖）。"""
    import sqlite3
    con = sqlite3.connect(path)
    try:
        cur = con.cursor()
        lines = []
        for (name,) in cur.execute(
                "SELECT name FROM sqlite_master WHERE type='table'"):
            try:
                rows = con.execute('SELECT * FROM "%s"' % name)
                cols = [d[0] for d in rows.description]
                lines.append("TABLE %s (%s)" % (name, ", ".join(cols)))
                for row in rows.fetchall():
                    lines.append(" | ".join("" if v is None else str(v) for v in row))
            except Exception:
                pass
        return "\n".join(lines)
    finally:
        con.close()


def _extract_text(ext, path):
    """返回 (文本, 缺失库名或None, 原因或None)。

    text=None 表示未成功抽取。原因 reason 取值：
      None      成功抽取到文本
      'no_lib'  缺少第三方抽取库（ImportError）
      'encrypted' 文件加密（异常信息含 password/encrypt/decrypt/protected）
      'image_only' PDF 抽到空文本（纯图片型/扫描件 PDF）
      'empty'    其他文档抽到空文本
      'error'    其他抽取异常（可能加密/损坏）
    """
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        lines.append(" ".join(cells))
            text = "\n".join(lines)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(r.text for r in para.runs)
                            if t:
                                lines.append(t)
            text = "\n".join(lines)
        elif ext == ".pdf":
            from pdfminer.high_level import extract_text as _pdf_text
            text = _pdf_text(path) or ""
        elif ext == ".db":
            text = _extract_sqlite(path)
        else:
            return (None, None, "error")
    except ImportError:
        return (None, _lib_for(ext), "no_lib")
    except Exception as e:
        msg = str(e).lower()
        if any(k in msg for k in ("password", "encrypt", "decrypt", "protected",
                                  "not a database")):
            return (None, None, "encrypted")
        # pdfminer 对真加密 PDF 常抛 PDFPasswordIncorrect（消息可能为空）；
        # 此时用 /Encrypt 标记二次确认，避免与“文件损坏”混淆
        if ext == ".pdf" and _pdf_has_encrypt(path):
            return (None, None, "encrypted")
        return (None, None, "error")
    # 抽到文本但为空：可能是纯图片型扫描件，也可能是加密 PDF
    # （pdfminer 对加密 PDF 常返回空而不报错，故额外探测 /Encrypt 标记）
    if not text or not text.strip():
        if ext == ".pdf" and _pdf_has_encrypt(path):
            return (None, None, "encrypted")
        return (None, None, "image_only" if ext == ".pdf" else "empty")
    return (text, None, None)


def _pdf_has_encrypt(path):
    """探测 PDF 是否含 /Encrypt 字典（粗略但高效，仅读前若干字节）。"""
    try:
        with open(path, "rb") as fb:
            head = fb.read(5 * 1024 * 1024)  # 至多 5MB，足够覆盖绝大多数文档
        return b"/Encrypt" in head
    except Exception:
        return False


def _process_extracted(path, mode, out_root, src_root, token_map, mapping,
                       names, counts_total, do_write, patterns, skipped):
    ext = os.path.splitext(path)[1].lower()
    rel_path = os.path.relpath(path, src_root)
    text, missing, reason = _extract_text(ext, path)
    if text is None:
        if missing:
            msg = ("缺少抽取库 %s，无法处理该二进制文档：请安装对应库"
                   "（docx→python-docx，xlsx→openpyxl，pptx→python-pptx，"
                   "pdf→pdfminer.six）后再处理" % missing)
        else:
            msg = REMIND.get(reason, "文本抽取失败（请确认文件未加密且未损坏）")
        if skipped is not None:
            skipped.append((path, msg))
        print("  [跳过] %s：%s" % (rel_path, msg), file=sys.stderr)
        return "skipped"
    if not do_write:
        c = scan_text(text, names, patterns)
        if c:
            print("  %-40s %s" % (rel_path, c))
            for k, v in c.items():
                counts_total[k] = counts_total.get(k, 0) + v
        return "processed"
    new_text, hits = desensitize_text(text, mode, token_map, counts_total, names, patterns)
    for h in hits:
        h["file"] = rel_path
        mapping.setdefault(rel_path, []).append(h)
    out_path = os.path.join(out_root, rel_path) + ".txt"
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(new_text)
    return "processed"


# ---------------------------------------------------------------------------
# 4. 映射表加密
# ---------------------------------------------------------------------------
def derive_key_from_passphrase(passphrase: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=200_000)
    return urlsafe_b64(kdf.derive(passphrase.encode()))


def urlsafe_b64(raw: bytes) -> bytes:
    import base64
    return base64.urlsafe_b64encode(raw)


def make_fernet(key: bytes) -> Fernet:
    return Fernet(key)


def encrypt_mapping(mapping: dict, keys_dir: str, passphrase: str = None) -> dict:
    """加密映射表。返回 {mapping_file, key_file|None, salt_file?}。"""
    os.makedirs(keys_dir, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    fernet = None
    key_file = None
    if passphrase:
        salt = secrets.token_bytes(16)
        key = derive_key_from_passphrase(passphrase, salt)
        fernet = make_fernet(key)
        salt_path = os.path.join(keys_dir, "salt_%s.bin" % ts)
        with open(salt_path, "wb") as f:
            f.write(salt)
    else:
        key = Fernet.generate_key()
        fernet = make_fernet(key)
        key_file = os.path.join(keys_dir, "desensitize_key_%s.key" % ts)
        with open(key_file, "wb") as f:
            f.write(key)
        os.chmod(key_file, 0o600)

    payload = json.dumps(mapping, ensure_ascii=False, indent=2).encode("utf-8")
    enc = fernet.encrypt(payload)
    enc_path = os.path.join(keys_dir, "mapping_%s.json.enc" % ts)
    with open(enc_path, "wb") as f:
        f.write(enc)
    os.chmod(enc_path, 0o600)
    return {"mapping_file": enc_path, "key_file": key_file, "encrypted": True}


def find_collisions(mapping: dict):
    """唯一性检测：找出"同一脱敏值 → 多个不同原始值"的多对一歧义。

    返回 {replacement: [去重后的原始值列表]}（仅含歧义项）。
    恢复（decrypt）时，若某 replacement 对应多个 original，则无法唯一还原，存在混淆风险。
    - token 模式：每个原始值分配唯一令牌，天然不会进入此表。
    - hybrid 模式：语义掩码 + 唯一令牌，replacement 含唯一令牌，天然不会进入此表。
    - mask 模式：中文姓名/地址/机构、邮箱等同型值易碰撞，需据此告警。
    - redact 模式：不可逆，恢复本就不适用。
    """
    rev = {}
    for items in mapping.values():
        for it in items:
            rev.setdefault(it["replacement"], set()).add(it["original"])
    return {r: sorted(v) for r, v in rev.items() if len(v) > 1}


def _report_skipped(skipped, src_root):
    """打印“未处理文件清单”（skip manifest），杜绝目录扫描的静默漏扫。"""
    if not skipped:
        return
    print("\n[警告] 以下 %d 个文件未被处理（可能含敏感信息，上云前须先转文本/OCR/"
          "安装抽取库，或人工复核）：" % len(skipped), file=sys.stderr)
    for p, reason in skipped:
        print("    - %s  （%s）" % (os.path.relpath(p, src_root), reason), file=sys.stderr)


def report_collisions(collisions: dict, mode: str):
    """打印碰撞告警（到 stderr），并给出模式建议。"""
    if not collisions:
        return
    print("\n[警告] 检测到 %d 个脱敏值存在“多对一”歧义（恢复时可能混淆）："
          % len(collisions), file=sys.stderr)
    for r, origs in collisions.items():
        preview = " / ".join(origs[:8]) + (" …" if len(origs) > 8 else "")
        print("    %s  <-  %s" % (r, preview), file=sys.stderr)
    if mode == "mask":
        print("  → 原因：mask 模式按类型生成固定掩码，同型不同值会碰撞（如不同姓名同为"
              "\"张*\"）。", file=sys.stderr)
        print("  → 建议：如需无歧义恢复，请改用 --mode hybrid（语义掩码+唯一令牌，"
              "保留字段语义且天然 1:1）或 --mode token。", file=sys.stderr)


# ---------------------------------------------------------------------------
# 5. 命令行
# ---------------------------------------------------------------------------
def iter_targets(input_path: str, recursive: bool):
    """产出待处理文件。**注意：这里产出目录下全部文件**，由 process_file 按扩展名
    决定处理或登记为“未处理”（skip manifest），从而杜绝“静默漏扫”。"""
    if os.path.isfile(input_path):
        yield input_path
        return
    if recursive:
        for root, _, files in os.walk(input_path):
            for fn in files:
                yield os.path.join(root, fn)
    else:
        for fn in sorted(os.listdir(input_path)):
            p = os.path.join(input_path, fn)
            if os.path.isfile(p):
                yield p


def load_names(path: str):
    if not path:
        return set()
    try:
        with open(path, "r", encoding="utf-8") as f:
            return {line.strip() for line in f if line.strip()}
    except FileNotFoundError:
        print("  [警告] 姓名清单文件不存在：%s（已忽略）" % path, file=sys.stderr)
        return set()


def cmd_scan(args):
    names = load_names(args.names)
    patterns = build_patterns(args.cn_enhance)
    input_path = os.path.abspath(args.input)
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(input_path)
    total = {}
    skipped = []
    print("扫描命中（不生成文件）%s：" % ("[中文增强开启]" if args.cn_enhance else ""))
    for p in iter_targets(input_path, args.recursive):
        process_file(p, "mask", "", src_root, {}, {}, names, total, do_write=False,
                     patterns=patterns, skipped=skipped)
    if total:
        print("\n汇总：", json.dumps(total, ensure_ascii=False))
    else:
        print("未发现已知敏感标识符。")
    _report_skipped(skipped, src_root)


def cmd_run(args):
    names = load_names(args.names)
    patterns = build_patterns(args.cn_enhance)
    input_path = os.path.abspath(args.input)
    src_root = input_path if os.path.isdir(input_path) else os.path.dirname(input_path)
    out_root = args.out
    keys_dir = args.keys
    os.makedirs(out_root, exist_ok=True)

    mapping = {}
    token_map = {}
    total = {}
    skipped = []
    print("开始脱敏（模式=%s）%s..." % (args.mode, "[中文增强开启]" if args.cn_enhance else ""))
    for p in iter_targets(input_path, args.recursive):
        process_file(p, args.mode, out_root, src_root, token_map, mapping, names, total,
                     do_write=True, patterns=patterns, skipped=skipped)

    enc_info = encrypt_mapping(mapping, keys_dir, args.passphrase)
    # 唯一性检测：仅 mask 模式可能因“同型不同值”产生多对一歧义（影响恢复）。
    # token / hybrid 由唯一令牌保证 1:1，天然无歧义；redact 不可逆，恢复本就不适用——
    # 二者都不应进入碰撞告警（否则 redact 会因 [REDACTED_x] 全相同而误报告警）。
    if args.mode == "mask":
        collisions = find_collisions(mapping)
    else:
        collisions = {}
    if args.mode == "redact":
        safety = "irreversible"
    elif collisions:
        safety = "ambiguous"
    else:
        safety = "unique"
    report_collisions(collisions, args.mode)
    # 去掉映射表里的原始值明细，仅保留统计写入报告（原始值已在加密文件中）
    report = {
        "time": datetime.now().isoformat(timespec="seconds"),
        "input": os.path.abspath(args.input),
        "out": os.path.abspath(out_root),
        "mode": args.mode,
        "counts": total,
        "items_encrypted": sum(len(v) for v in mapping.values()),
        "restoration_safety": safety,
        "collision_risks": collisions,
        "mapping_encrypted": enc_info["mapping_file"],
        "key_file": enc_info.get("key_file"),
        "skipped_files": [os.path.relpath(p, src_root) for p, _ in skipped],
        "note": "原始敏感值仅存于加密映射表；脱敏副本位于 out 目录，可上云。"
                "restoration_safety=unique 表示可无歧义恢复（token / hybrid 模式恒为 unique）；"
                "=ambiguous 仅可能出现在 mask 模式（同型不同值碰撞），恢复可能混淆，"
                "建议改用 --mode hybrid（语义掩码+唯一令牌，兼顾字段语义与无歧义）或 --mode token；"
                "=irreversible 表示 redact 模式（不可逆，不可还原）。",
    }
    report_path = os.path.join(out_root, "..", "desensitize_report.json")
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    print("\n完成。")
    print("  脱敏副本目录 : %s" % os.path.abspath(out_root))
    print("  加密映射表   : %s" % enc_info["mapping_file"])
    if enc_info.get("key_file"):
        print("  密钥文件     : %s  (权限 600，请勿与脱敏副本一同上传)" % enc_info["key_file"])
    else:
        print("  密钥         : 由 --passphrase 派生（salt 已存于 keys 目录）")
    print("  命中统计     : %s" % json.dumps(total, ensure_ascii=False))
    print("  恢复安全性   : %s" % safety +
          ("（存在 %d 处多对一碰撞，恢复可能混淆）" % len(collisions)
           if safety == "ambiguous" else ""))
    print("  报告         : %s" % os.path.abspath(report_path))
    _report_skipped(skipped, src_root)


def _load_mapping(keys_dir, file=None, key=None, passphrase=None):
    """解密并返回映射表（data, enc_path）。供 decrypt / restore 复用。"""
    if file:
        enc_path = file
    else:
        encs = sorted(glob.glob(os.path.join(keys_dir, "mapping_*.json.enc")))
        if not encs:
            raise FileNotFoundError("未在 %s 找到 mapping_*.json.enc" % keys_dir)
        enc_path = encs[-1]
    if passphrase:
        salts = sorted(glob.glob(os.path.join(keys_dir, "salt_*.bin")))
        if not salts:
            raise ValueError("未找到 salt 文件，无法用口令解密")
        salt = open(salts[-1], "rb").read()
        k = derive_key_from_passphrase(passphrase, salt)
    else:
        keys = sorted(glob.glob(os.path.join(keys_dir, "desensitize_key_*.key")))
        if not keys:
            raise ValueError("未在 %s 找到 desensitize_key_*.key（也未提供 --passphrase）"
                             % keys_dir)
        k = open(key or keys[-1], "rb").read()
    fernet = make_fernet(k)
    data = json.loads(fernet.decrypt(open(enc_path, "rb").read()).decode("utf-8"))
    return data, enc_path


def cmd_decrypt(args):
    """解密映射表，供本地复核可逆性（无需上云）。"""
    try:
        data, enc_path = _load_mapping(args.keys, args.file, args.key, args.passphrase)
    except (FileNotFoundError, ValueError) as e:
        print(str(e), file=sys.stderr)
        return
    except Exception as e:
        print("解密失败：%s" % e, file=sys.stderr)
        return

    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print("已解密映射表写入：%s" % os.path.abspath(args.out))
        return

    total = sum(len(v) for v in data.values())
    # 过滤 redact 模式产生的 [REDACTED_*]：不可逆、本不适用恢复，不应作为歧义告警
    collisions = {r: v for r, v in find_collisions(data).items()
                  if not str(r).startswith("[REDACTED_")}
    print("映射表文件：%s" % enc_path)
    print("涉及文件 %d 个，命中条目 %d 条：" % (len(data), total))
    for f, items in data.items():
        print("  %-40s %d 条" % (f, len(items)))
    if collisions:
        print("\n[提醒] 该映射表存在 %d 处“同一脱敏值→多个原始值”的歧义，恢复时可能混淆："
              % len(collisions), file=sys.stderr)
        for r, origs in collisions.items():
            preview = " / ".join(origs[:8]) + (" …" if len(origs) > 8 else "")
            print("    %s  <-  %s" % (r, preview), file=sys.stderr)
        print("  → 如需无歧义恢复，请对该数据改用 --mode hybrid 或 --mode token 重新脱敏。",
              file=sys.stderr)
    else:
        print("\n[OK] 未检测到多对一歧义，原始值可唯一还原"
              "（hybrid / token 模式，或碰撞为空的 mask 模式）。")
    print("\n（加 --out <文件.json> 可导出完整 original↔replacement 明细，用于本地复核）")


def cmd_restore(args):
    """用加密映射表把脱敏副本回填为含原值的内部文档（映射表不离本地）。"""
    try:
        data, enc_path = _load_mapping(args.keys, args.file, args.key, args.passphrase)
    except (FileNotFoundError, ValueError) as e:
        print(str(e), file=sys.stderr)
        return
    except Exception as e:
        print("解密失败：%s" % e, file=sys.stderr)
        return

    input_root = args.input
    out_root = args.out
    os.makedirs(out_root, exist_ok=True)
    types = set(t.strip() for t in (args.types or "").split(",") if t.strip())

    total_restored = 0
    total_collision = 0
    total_skipped = 0
    for rel_path, items in data.items():
        # 定位脱敏副本：文本/代码按原扩展名；抽取型(.db/.pdf…)实际输出为 .txt
        cand = os.path.join(input_root, rel_path)
        suffix = ""
        if not os.path.isfile(cand):
            cand2 = cand + ".txt"
            if os.path.isfile(cand2):
                cand = cand2
                suffix = ".txt"
            else:
                print("  [跳过] 未找到脱敏副本：%s" % rel_path, file=sys.stderr)
                total_skipped += 1
                continue
        with open(cand, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        # 构建 replacement -> original；检测多对一歧义（mask 同型/*** 密钥）
        repl_map = {}
        collisions = {}
        for it in items:
            if types and it.get("type") not in types:
                continue
            r = it["replacement"]
            o = it["original"]
            if r in repl_map and repl_map[r] != o:
                collisions.setdefault(r, set([repl_map[r]])).add(o)
            else:
                repl_map[r] = o
        for r in collisions:
            repl_map.pop(r, None)
        if collisions:
            total_collision += len(collisions)
            for r, origs in collisions.items():
                preview = " / ".join(sorted(origs))
                print("  [警告] 替换串 %r 对应多个原始值（%s），无法唯一还原，已跳过"
                      % (r, preview[:120]), file=sys.stderr)

        # 单次遍历替换（最长优先，避免子串互相干扰）；令牌模式请确保数据不含同名串
        if repl_map:
            escaped = sorted(((re.escape(r), r) for r in repl_map),
                             key=lambda x: -len(x[0]))
            pat = re.compile("|".join(e[0] for e in escaped))
            content = pat.sub(lambda m: repl_map[m.group(0)], content)
            total_restored += len(repl_map)

        out_path = os.path.join(out_root, rel_path) + suffix
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(content)
        print("  回填 -> %s （%d 处）" % (os.path.relpath(out_path, out_root), len(repl_map)))

    print("\n完成。回填目录: %s" % os.path.abspath(out_root))
    print("  回填条目: %d | 因碰撞跳过: %d | 未找到副本: %d" % (
        total_restored, total_collision, total_skipped))
    print("  （映射表 %s 已用于回填；回填产物恢复为含原值的本地内部文档，请按需谨慎保管，"
          "勿随脱敏副本一同外传）" % enc_path)


def cmd_audit(args):
    """基于 run 生成的 desensitize_report.json 自动产出 11 项审计文档。"""
    report_path = args.report
    if not os.path.isfile(report_path):
        print("未找到报告文件：%s" % report_path, file=sys.stderr)
        return
    rep = json.load(open(report_path, encoding="utf-8"))
    time = rep.get("time", "")
    inp = rep.get("input", "")
    out = rep.get("out", "")
    mode = rep.get("mode", "")
    counts = rep.get("counts", {})
    safety = rep.get("restoration_safety", "")
    collisions = rep.get("collision_risks", {})
    mapping_file = rep.get("mapping_encrypted", "")
    key_file = rep.get("key_file") or "（由 --passphrase 派生，salt 存于 keys 目录）"
    skipped = rep.get("skipped_files", [])
    items = rep.get("items_encrypted", 0)

    # 简化分级：直接标识符 / 密钥类 → 高；其余 → 中
    HIGH = {"id_card", "phone", "bank_card", "email", "ip", "plate",
            "passport", "name", "cn_name", "secret"}
    levels = {}
    for k, v in counts.items():
        lvl = "高" if k in HIGH else "中"
        levels.setdefault(lvl, {})[k] = v

    L = []
    L.append("# 脱敏审计记录（由 desensitize.py audit 自动生成）")
    L.append("")
    L.append("- 生成时间：%s" % datetime.now().isoformat(timespec="seconds"))
    L.append("- 任务时间：%s" % time)
    L.append("- 任务描述：基于 `run` 报告自动生成；业务背景请人工补充")
    L.append("")
    L.append("## 一、数据盘点")
    L.append("- 原始数据：%s" % inp)
    L.append("- 脱敏条目（已加密入映射表）：%s 条" % items)
    L.append("- 脱敏副本目录：%s" % out)
    L.append("")
    L.append("## 二、敏感字段与级别")
    if levels:
        for lvl in ("高", "中", "低"):
            if lvl in levels:
                detail = "，".join("%s×%d" % (k, v)
                                  for k, v in sorted(levels[lvl].items()))
                L.append("- %s风险：%s" % (lvl, detail))
    else:
        L.append("- （报告未记录命中，或仅 redact 模式）")
    L.append("")
    L.append("## 三、脱敏方法")
    L.append("- 模式：%s" % mode)
    L.append("- 恢复安全性：%s%s" % (
        safety, "（存在碰撞，恢复可能混淆，建议改用 hybrid）" if safety == "ambiguous" else ""))
    if mode == "redact":
        L.append("- redact 不可逆，被抑制字段不可还原（精度：相关字段已删除）")
    else:
        L.append("- 数值精度：仅去标识、保留精确数值（未泛化/随机化），满足财务/审计红线")
    L.append("")
    L.append("## 四、密钥与映射表（重识别钥匙，须分离/加密/最小权限）")
    L.append("- 加密映射表：%s" % mapping_file)
    L.append("- 密钥文件：%s （权限 600）" % key_file)
    L.append("- 状态：原始文件与映射表留本地、未随副本上云（对应 SKILL.md 清单 1/4/5）")
    L.append("")
    L.append("## 五、上云内容")
    L.append("- 仅脱敏副本：%s" % out)
    L.append("")
    L.append("## 六、本地复核结论")
    if safety == "unique":
        L.append("- 令牌唯一、无歧义恢复：PASS（可用 `decrypt`/`restore` 还原且无混淆）")
    elif safety == "ambiguous":
        L.append("- 存在 %d 处多对一歧义，恢复可能混淆：" % len(collisions))
        for r, origs in collisions.items():
            L.append("    - %s <- %s" % (r, " / ".join(origs[:8])))
    else:
        L.append("- redact 不可逆，不适用恢复校验")
    L.append("")
    L.append("## 七、未处理文件（skip manifest，须人工复核）")
    if skipped:
        for s in skipped:
            L.append("- %s" % s)
        L.append("> 上述文件未脱敏，上云前须转文本/OCR/解密或确认不含敏感信息（SKILL.md 红线）")
    else:
        L.append("- 无（全部目标文件已处理或被显式忽略）")
    L.append("")
    L.append("## 八、外泄风险自评")
    risk = "低" if (safety == "unique" and not skipped) else "中"
    L.append("- 风险等级：%s" % risk)
    L.append("- 理由：脱敏副本可上云；原始/映射表留本地分离；%s"
             % ("存在歧义/未处理文件需关注" if risk == "中"
                else "无歧义恢复、无遗漏"))
    L.append("")
    L.append("## 九、操作人与待办")
    L.append("- 操作人：AI Agent（本地执行）")
    L.append("- 异常与待办：%s" % ("无" if not (collisions or skipped)
                                    else "见第六/七节，需人工复核后上云"))
    L.append("")
    L.append("> 本审计由 `desensitize.py audit` 自动生成；SKILL.md 规定的 11 项上云前自查"
             "清单与人工复核仍须由操作人逐项确认，禁止“一键脱敏即上云”。")

    out_md = args.out
    with open(out_md, "w", encoding="utf-8") as f:
        f.write("\n".join(L) + "\n")
    print("已生成审计文档：%s" % os.path.abspath(out_md))


def build_parser():
    p = argparse.ArgumentParser(description="一键脱敏本地脚本（数据不出本机）")
    sub = p.add_subparsers(dest="cmd", required=True)

    sp = sub.add_parser("scan", help="扫描并报告敏感字段命中（不生成文件）")
    sp.add_argument("input", help="文件或目录")
    sp.add_argument("--recursive", action="store_true", help="递归处理目录")
    sp.add_argument("--names", help="已知姓名清单文件（每行一个）")
    sp.add_argument("--cn-enhance", action="store_true",
                    help="中文识别增强：额外识别中文姓名/地址/机构名（本地正则，离线）")
    sp.set_defaults(func=cmd_scan)

    rp = sub.add_parser("run", help="脱敏并生成加密映射表")
    rp.add_argument("input", help="文件或目录")
    rp.add_argument("--out", default="./desensitized", help="脱敏副本输出目录（默认 ./desensitized）")
    rp.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录（默认 ./.desensitize_keys）")
    rp.add_argument("--mode", choices=["mask", "token", "redact", "hybrid"], default="hybrid",
                    help="脱敏模式（默认 hybrid）：mask=掩码(可逆) / token=令牌化(可逆,跨记录关联) / "
                         "hybrid=语义掩码+唯一令牌(可逆,无歧义,保留字段语义) / redact=抑制(不可逆)")
    rp.add_argument("--passphrase", default=None, help="用口令派生密钥（密钥不落盘）")
    rp.add_argument("--recursive", action="store_true", help="递归处理目录")
    rp.add_argument("--names", help="已知姓名清单文件（每行一个）")
    rp.add_argument("--cn-enhance", action="store_true",
                    help="中文识别增强：额外识别中文姓名/地址/机构名（本地正则，离线）")
    rp.set_defaults(func=cmd_run)

    dp = sub.add_parser("decrypt", help="解密映射表，供本地复核可逆性（不依赖上云）")
    dp.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录")
    dp.add_argument("--file", default=None, help="指定 mapping_*.json.enc（默认取目录中最新一个）")
    dp.add_argument("--key", default=None, help="指定密钥文件（默认取目录中最新 .key）")
    dp.add_argument("--passphrase", default=None, help="若用口令派生密钥，提供同一口令")
    dp.add_argument("--out", default=None, help="导出解密后的映射表 JSON 到该路径")
    dp.set_defaults(func=cmd_decrypt)

    rp2 = sub.add_parser("restore", help="用映射表把脱敏副本回填为含原值的内部文档")
    rp2.add_argument("--keys", default="./.desensitize_keys", help="加密映射表目录")
    rp2.add_argument("--input", default="./desensitized",
                     help="脱敏副本目录/文件（默认 ./desensitized）")
    rp2.add_argument("--out", default="./restored",
                     help="回填后内部文档输出目录（默认 ./restored）")
    rp2.add_argument("--file", default=None, help="指定 mapping_*.json.enc（默认取目录中最新一个）")
    rp2.add_argument("--key", default=None, help="指定密钥文件（默认取目录中最新 .key）")
    rp2.add_argument("--passphrase", default=None, help="若用口令派生密钥，提供同一口令")
    rp2.add_argument("--types", default=None,
                     help="仅回填指定类型，逗号分隔（如 name,id_card）；默认全部回填")
    rp2.set_defaults(func=cmd_restore)

    ap = sub.add_parser("audit", help="基于 run 报告自动生成 11 项脱敏审计文档")
    ap.add_argument("--report", default="./desensitize_report.json",
                    help="run 生成的报告（默认 ./desensitize_report.json）")
    ap.add_argument("--out", default="./desensitize_audit.md",
                    help="审计文档输出路径（默认 ./desensitize_audit.md）")
    ap.set_defaults(func=cmd_audit)
    return p


def main():
    args = build_parser().parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
